#!/usr/bin/env python3
"""Generate ZTO-Pechoin warehouse logistics proposal PPT — v2 redesigned."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ZTO_LOGO = os.path.join(SCRIPT_DIR, "zto_logo.png")
PECHOIN_LOGO = os.path.join(SCRIPT_DIR, "pechoin_logo.png")
OUTPUT = os.path.join(SCRIPT_DIR, "中通-百雀羚仓储物流方案.pptx")

# ── Color Palette ──
PRIMARY = RGBColor(0x1B, 0x5E, 0xA0)       # deep professional blue
PRIMARY_LIGHT = RGBColor(0xE8, 0xF0, 0xFA)  # very light blue bg
ACCENT = RGBColor(0x2E, 0x8B, 0x57)         # sea green (Pechoin)
ACCENT_LIGHT = RGBColor(0xE8, 0xF6, 0xEF)   # very light green bg
DARK_BG = RGBColor(0x0D, 0x2B, 0x45)        # near-navy for dark slides
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x2D, 0x2D, 0x2D)
DARK_GRAY = RGBColor(0x55, 0x55, 0x55)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF4)
BORDER_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED = RGBColor(0xC0, 0x39, 0x2B)
CARD_BG = RGBColor(0xFA, 0xFA, 0xFA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.7)
CONTENT_W = SLIDE_W - MARGIN * 2

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg_color(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def shape(slide, st, l, t, w, h, fill=None, line=None, lw=1):
    s = slide.shapes.add_shape(st, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def txt(tf_or_shape, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tf = tf_or_shape.text_frame if hasattr(tf_or_shape, 'text_frame') else tf_or_shape
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.alignment = align
    return p


def para(tf, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT, sb=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.alignment = align
    p.space_before = sb
    return p


def tbox(slide, l, t, w, h, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.alignment = align
    return box


def page_header(slide, title, subtitle=None):
    """Clean minimal header: thin accent line + title."""
    shape(slide, MSO_SHAPE.RECTANGLE, MARGIN, Inches(0.45), Inches(0.35), Pt(4), fill=PRIMARY)
    tbox(slide, MARGIN + Inches(0.5), Inches(0.3), Inches(8), Inches(0.5),
         title, size=26, bold=True, color=PRIMARY)
    if subtitle:
        tbox(slide, MARGIN + Inches(0.5), Inches(0.75), Inches(10), Inches(0.35),
             subtitle, size=12, color=MID_GRAY)
    # subtle bottom line
    shape(slide, MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.15), CONTENT_W, Pt(0.75), fill=BORDER_GRAY)


def page_footer(slide, num):
    tbox(slide, SLIDE_W - Inches(1.2), Inches(7.05), Inches(0.8), Inches(0.3),
         str(num), size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)


def card(slide, l, t, w, h, fill_color=CARD_BG, border_color=BORDER_GRAY):
    return shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, fill=fill_color, line=border_color, lw=0.75)


def add_img(slide, path, l, t, w=None, h=None):
    if os.path.exists(path):
        return slide.shapes.add_picture(path, l, t, w, h)
    return None


def arrow_right(slide, x1, y1, x2, y2, color=PRIMARY):
    """Clean horizontal arrow."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(1.5)
    ln = conn.line._get_or_add_ln()
    for ch in list(ln):
        if 'headEnd' in ch.tag or 'tailEnd' in ch.tag:
            ln.remove(ch)
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return conn


def arrow_down(slide, x, y1, y2, color=PRIMARY):
    """Clean vertical arrow."""
    return arrow_right(slide, x, y1, x, y2, color)


def flow_box(slide, l, t, w, h, text, fill=PRIMARY, size=11, tc=WHITE, bold=False):
    s = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, fill=fill)
    tf = s.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = tc
    p.font.name = "Microsoft YaHei"
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return s


def styled_table(slide, l, t, w, h, rows, cols):
    ts = slide.shapes.add_table(rows, cols, l, t, w, h)
    return ts.table


def set_cell(cell, text, size=10, bold=False, color=BLACK, fill=None, align=PP_ALIGN.CENTER):
    cell.text = text
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.alignment = align
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


# ══════════════════════════════════════════════
#  S1: Cover
# ══════════════════════════════════════════════
s = blank_slide()
bg_color(s, WHITE)

# Left accent strip
shape(s, MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.06), SLIDE_H, fill=PRIMARY)

# Logos row — centered, proper sizing
logo_y = Inches(1.8)
# ZTO logo
add_img(s, ZTO_LOGO, Inches(4.8), logo_y, w=Inches(1.2))
# "×" separator
tbox(s, Inches(6.1), logo_y + Inches(0.25), Inches(0.8), Inches(0.6),
     "×", size=36, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)
# Pechoin logo on green pill
pechoin_pill = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), logo_y + Inches(0.1), Inches(1.8), Inches(0.85), fill=ACCENT)
add_img(s, PECHOIN_LOGO, Inches(7.2), logo_y + Inches(0.22), w=Inches(1.4))

# Title
tbox(s, Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.7),
     "百雀羚工厂调拨及电商仓配送", size=32, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
tbox(s, Inches(1.5), Inches(3.85), Inches(10.3), Inches(0.7),
     "仓储物流运输方案", size=40, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

# Divider line
shape(s, MSO_SHAPE.RECTANGLE, Inches(5.4), Inches(4.75), Inches(2.5), Pt(2), fill=ACCENT)

# Subtitle
tbox(s, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.4),
     "整车运输  |  厂间调拨  |  电商仓直送", size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Bottom info
tbox(s, Inches(1.5), Inches(5.9), Inches(10.3), Inches(0.35),
     "2026年7月", size=13, color=MID_GRAY, align=PP_ALIGN.CENTER)
tbox(s, Inches(1.5), Inches(6.35), Inches(10.3), Inches(0.4),
     "中通快递股份有限公司", size=15, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
#  S2: 项目概述
# ══════════════════════════════════════════════
s = blank_slide()
bg_color(s)
page_header(s, "项目概述", "百雀羚工厂调拨纯运输项目 · 整车运输服务")

# Left: 项目背景
bg_card = card(s, MARGIN, Inches(1.5), Inches(5.6), Inches(2.4), fill_color=LIGHT_GRAY)
tf = bg_card.text_frame
tf.word_wrap = True
txt(tf, "项目背景", 16, True, PRIMARY)
para(tf, "百雀羚是国内知名护肤品牌，产品覆盖全线护肤、彩妆及个护品类。", 12)
para(tf, "本项目为其工厂间调拨及电商仓配送提供整车运输服务，", 12)
para(tf, "涵盖三大生产基地与电商大仓之间的整车物流需求。", 12)

# Left: 运输范围
range_card = card(s, MARGIN, Inches(4.1), Inches(5.6), Inches(2.7), fill_color=LIGHT_GRAY)
tf = range_card.text_frame
tf.word_wrap = True
txt(tf, "运输范围", 16, True, PRIMARY)
para(tf, "业务一：厂间调拨", 13, True, BLACK, sb=Pt(10))
para(tf, "三分厂 / 四分厂 / 菲尼克斯工厂之间整车调拨，承运商负责装卸", 11, color=DARK_GRAY, sb=Pt(2))
para(tf, "业务二：电商仓直送", 13, True, BLACK, sb=Pt(10))
para(tf, "三个工厂发往电商大仓的整车业务，电商仓端承运商不需要装卸服务", 11, color=DARK_GRAY, sb=Pt(2))

# Right: 关键数据卡片 (2×2 grid)
tbox(s, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.4),
     "项目关键数据", size=18, bold=True, color=DARK_BG)

metrics = [
    ("3", "生产基地", "三分厂 · 四分厂 · 菲尼克斯"),
    ("2", "业务类型", "厂间调拨 + 电商仓配送"),
    ("7×24", "运营保障", "全天候客服响应机制"),
    ("98%+", "准时交付率", "全流程跟踪与品质保障"),
]
for i, (num, title, desc) in enumerate(metrics):
    cx = Inches(6.8) + Inches(2.9) * (i % 2)
    cy = Inches(2.1) + Inches(1.65) * (i // 2)
    c = card(s, cx, cy, Inches(2.7), Inches(1.45), fill_color=WHITE, border_color=PRIMARY)
    tf = c.text_frame
    tf.word_wrap = True
    txt(tf, num, 26, True, PRIMARY, PP_ALIGN.CENTER)
    para(tf, title, 12, True, BLACK, PP_ALIGN.CENTER, sb=Pt(2))
    para(tf, desc, 9, False, MID_GRAY, PP_ALIGN.CENTER, sb=Pt(2))

page_footer(s, 2)


# ══════════════════════════════════════════════
#  S3: 运输网络
# ══════════════════════════════════════════════
s = blank_slide()
bg_color(s)
page_header(s, "运输网络", "三大生产基地调拨与电商仓配送路线")

# Layout: left column = 3 factories, right column = 电商大仓
# Factories vertically stacked, arrows to a central convergence point, then to 电商大仓

fac_x = MARGIN
fac_w = Inches(2.6)
fac_h = Inches(0.75)
fac_gap = Inches(0.6)
fac_y0 = Inches(1.7)

factories = ["三分厂", "四分厂", "菲尼克斯工厂"]
fac_boxes = []
for i, name in enumerate(factories):
    fy = fac_y0 + (fac_h + fac_gap) * i
    b = flow_box(s, fac_x, fy, fac_w, fac_h, name, fill=ACCENT, size=14, bold=True)
    fac_boxes.append(b)

# 电商大仓 — vertically centered
ec_x = Inches(9.5)
ec_y = Inches(2.55)
ec_w = Inches(2.8)
ec_h = Inches(1.1)
flow_box(s, ec_x, ec_y, ec_w, ec_h, "电商大仓", fill=PRIMARY, size=18, bold=True)

# Draw clean horizontal arrows from each factory to a merge zone, then to 电商大仓
merge_x = Inches(6.5)  # merge zone x position
fac_right = fac_x + fac_w

for i in range(3):
    fy = fac_y0 + (fac_h + fac_gap) * i + fac_h // 2
    # Horizontal line from factory to merge zone
    arrow_right(s, fac_right + Inches(0.1), fy, merge_x, fy, ACCENT)

# Merge lines to 电商大仓
ec_center_y = ec_y + ec_h // 2
for i in range(3):
    fy = fac_y0 + (fac_h + fac_gap) * i + fac_h // 2
    # Diagonal from merge zone to ec
    arrow_right(s, merge_x, fy, ec_x - Inches(0.1), ec_center_y, PRIMARY)

# Also draw factory-to-factory arrows (bidirectional, on left side)
# Use simple horizontal lines between adjacent factories
for i in range(2):
    fy1 = fac_y0 + (fac_h + fac_gap) * i + fac_h // 2
    fy2 = fac_y0 + (fac_h + fac_gap) * (i + 1) + fac_h // 2
    mid_fy = (fy1 + fy2) // 2
    # Small arrow indicator
    arrow_down(s, fac_x + fac_w + Inches(0.05), fy1, fy2, ACCENT)
    arrow_right(s, fac_x + fac_w + Inches(0.05), fy2, fac_x + fac_w + Inches(0.05), fy1, ACCENT)

# Labels
lbl1 = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(1.3), Inches(3.0), Inches(0.42),
             fill=ACCENT_LIGHT, line=ACCENT, lw=0.75)
tf = lbl1.text_frame
txt(tf, "厂间调拨（承运商装卸）", 10, False, ACCENT, PP_ALIGN.CENTER)
lbl1.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

lbl2 = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(4.3), Inches(3.0), Inches(0.42),
             fill=PRIMARY_LIGHT, line=PRIMARY, lw=0.75)
tf = lbl2.text_frame
txt(tf, "电商仓直送（无需装卸）", 10, False, PRIMARY, PP_ALIGN.CENTER)
lbl2.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Bottom info card
info = card(s, MARGIN, Inches(5.5), CONTENT_W, Inches(1.5), fill_color=LIGHT_GRAY)
tf = info.text_frame
tf.word_wrap = True
txt(tf, "运输说明", 14, True, DARK_BG)
para(tf, "厂间调拨：三分厂 / 四分厂 / 菲尼克斯工厂之间调拨，承运商负责装卸", 11)
para(tf, "电商仓配送：三厂发往电商大仓，承运商无需在电商仓端提供装卸服务", 11)
para(tf, "车辆类型：整车运输，根据货量灵活调配车型", 11)

page_footer(s, 3)


# ══════════════════════════════════════════════
#  S4: 业务流程图
# ══════════════════════════════════════════════
s = blank_slide()
bg_color(s)
page_header(s, "业务流程", "从货量下达到目的地交接的全流程管理")

# Horizontal flow — 4 steps in row 1, 4 steps in row 2
# Row 1: left → right, Row 2: right → left (snake layout)

step_w = Inches(2.55)
step_h = Inches(0.85)
step_gap = Inches(0.45)
x0 = MARGIN

# Row 1
r1_steps = [
    ("工厂提前一天\n下达货量任务", PRIMARY),
    ("中通调度中心\n接收并确认", PRIMARY),
    ("车辆调配\n司机指派", PRIMARY),
    ("车辆到达\n工厂装货", ACCENT),
]
y1 = Inches(1.6)
for i, (text, color) in enumerate(r1_steps):
    x = x0 + (step_w + step_gap) * i
    flow_box(s, x, y1, step_w, step_h, text, fill=color, size=11, bold=True)
    if i < 3:
        ax = x + step_w
        arrow_right(s, ax + Inches(0.03), y1 + step_h // 2, ax + step_gap - Inches(0.03), y1 + step_h // 2)

# Down arrow from last of row 1
last_x = x0 + (step_w + step_gap) * 3 + step_w // 2
arrow_down(s, last_x, y1 + step_h, Inches(3.0))

# Row 2 (right to left)
r2_steps = [
    ("回单签收\n流程闭环", ACCENT),
    ("到达目的地\n清点交接", ORANGE),
    ("运输途中\nGPS实时监控", PRIMARY),
    ("货物装车\n缠绕膜固定", ACCENT),
]
y2 = Inches(3.15)
for i, (text, color) in enumerate(r2_steps):
    x = x0 + (step_w + step_gap) * i
    flow_box(s, x, y2, step_w, step_h, text, fill=color, size=11, bold=True)
    if i < 3:
        ax = x + step_w
        arrow_right(s, ax + step_gap - Inches(0.03), y2 + step_h // 2, ax + Inches(0.03), y2 + step_h // 2)

# Timeline
tl_y = Inches(4.5)
tbox(s, MARGIN, tl_y - Inches(0.15), Inches(2), Inches(0.3),
     "时间节点：", 12, True, DARK_BG)
time_items = [
    ("T-1天", "工厂下达任务"),
    ("T-1天", "调度确认派车"),
    ("T日", "到厂装货"),
    ("T日", "运输配送"),
    ("T日", "到达交接"),
]
for i, (t, desc) in enumerate(time_items):
    x = Inches(2.0) + Inches(2.2) * i
    tbox(s, x, tl_y, Inches(1.5), Inches(0.28), t, 11, True, PRIMARY, PP_ALIGN.CENTER)
    tbox(s, x, tl_y + Inches(0.28), Inches(1.5), Inches(0.25), desc, 9, False, MID_GRAY, PP_ALIGN.CENTER)

# Timeline bar
shape(s, MSO_SHAPE.RECTANGLE, Inches(2.0), tl_y - Inches(0.02), Inches(9.5), Pt(1.5), fill=BORDER_GRAY)

# Key control points
kp = card(s, MARGIN, Inches(5.3), CONTENT_W, Inches(1.7), fill_color=LIGHT_GRAY)
tf = kp.text_frame
tf.word_wrap = True
txt(tf, "关键控制点", 14, True, DARK_BG)
para(tf, "① 工厂提前一天下达货量任务，中通调度中心15分钟内确认并反馈车辆安排", 11)
para(tf, "② 装货时司机核对货物数量、形态（整托/半托/散箱）并拍照留档", 11)
para(tf, "③ 运输全程GPS监控，异常停留/偏航自动预警，客服实时跟进", 11)
para(tf, "④ 到达目的地后双方清点交接，签收回单拍照上传，当日完成闭环", 11)

page_footer(s, 4)


# ══════════════════════════════════════════════
#  S5: 货物形态与装卸规范
# ══════════════════════════════════════════════
s = blank_slide()
bg_color(s)
page_header(s, "货物形态与装卸规范", "整托 / 半托 / 散箱 分类管理与操作标准")

col_w = Inches(3.7)
col_gap = Inches(0.35)
col_y = Inches(1.5)

cargo = [
    ("整托", "大部分", ACCENT, [
        "工厂出厂已缠绕好缠绕膜",
        "整托装卸，叉车作业",
        "厂间调拨 & 电商仓均适用",
        "装卸效率最高，破损率最低",
        "占比：约 70% ~ 80%",
    ]),
    ("半托", "部分", ORANGE, [
        "工厂缠绕好膜后装车",
        "需注意堆叠稳定性",
        "厂间调拨 & 电商仓均适用",
        "装卸时需加固防倾斜",
        "占比：约 15% ~ 25%",
    ]),
    ("散箱", "极少", RED, [
        "仅出现在送电商仓业务中",
        "需人工逐箱搬运装卸",
        "需额外防护与固定措施",
        "重点防止破损与丢失",
        "占比：约 5% 以内",
    ]),
]

for i, (title, sub, color, items) in enumerate(cargo):
    x = MARGIN + (col_w + col_gap) * i
    # Color header
    hdr = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, col_y, col_w, Inches(0.6), fill=color)
    tf = hdr.text_frame
    txt(tf, f"{title}（{sub}）", 15, True, WHITE, PP_ALIGN.CENTER)
    hdr.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Body
    body = card(s, x, col_y + Inches(0.68), col_w, Inches(2.3), fill_color=WHITE, border_color=color)
    tf = body.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    for j, item in enumerate(items):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = BLACK
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(7)

# Bottom note
note = card(s, MARGIN, Inches(4.85), CONTENT_W, Inches(1.85), fill_color=RGBColor(0xFF, 0xF8, 0xE1), border_color=ORANGE)
tf = note.text_frame
tf.word_wrap = True
txt(tf, "装卸服务说明", 14, True, ORANGE)
para(tf, "", 4)
para(tf, "厂间调拨：承运商负责装卸，配备专业装卸团队与叉车设备", 12)
para(tf, "电商仓直送：电商仓端承运商不需要装卸服务，由电商仓负责卸货", 12)
para(tf, "所有装卸过程拍照留档，异常破损即时上报并启动理赔流程", 12)

page_footer(s, 5)


# ══════════════════════════════════════════════
#  S6: 车辆管控方案
# ══════════════════════════════════════════════
s = blank_slide()
bg_color(s)
page_header(s, "车辆管控方案", "自有车辆与长期签约车辆的统一管理体系")

half_w = Emu(int((CONTENT_W - Inches(0.4)) / 2))

# Left: 自有车辆
lc = card(s, MARGIN, Inches(1.5), half_w, Inches(5.3), fill_color=PRIMARY_LIGHT, border_color=PRIMARY)
tf = lc.text_frame
tf.word_wrap = True
txt(tf, "自有车辆管控", 18, True, PRIMARY)
sections_l = [
    ("车辆档案管理", "一车一档：行驶证、营运证、保险、年检记录"),
    ("智能调度系统", "GPS定位 + 车载OBD，实时监控位置与车况"),
    ("定期维保计划", "按里程/时间双维度制定保养计划，确保车况良好"),
    ("安全检查", "出车前 / 行车中 / 收车后三阶段安全检查"),
    ("备用车辆", "配备应急备用车辆，保障运力不中断"),
]
for title, desc in sections_l:
    para(tf, title, 13, True, DARK_BG, sb=Pt(12))
    para(tf, desc, 11, color=DARK_GRAY, sb=Pt(2))

# Right: 签约车辆
rc = card(s, MARGIN + half_w + Inches(0.4), Inches(1.5), half_w, Inches(5.3), fill_color=ACCENT_LIGHT, border_color=ACCENT)
tf = rc.text_frame
tf.word_wrap = True
txt(tf, "长期签约车辆管控", 18, True, ACCENT)
sections_r = [
    ("准入审核", "资质审查：营业执照、道路运输许可证、车辆保险"),
    ("合同约束", "签订长期合作协议，明确服务标准与违约责任"),
    ("统一标准", "与自有车辆执行相同的装载、运输、交付标准"),
    ("绩效考核", "按月评估准时率、破损率、投诉率，末位淘汰"),
    ("运力储备", "保持20%以上的运力冗余，应对旺季和突发需求"),
]
for title, desc in sections_r:
    para(tf, title, 13, True, DARK_BG, sb=Pt(12))
    para(tf, desc, 11, color=DARK_GRAY, sb=Pt(2))

page_footer(s, 6)


# ══════════════════════════════════════════════
#  S7: 司机考核体系
# ══════════════════════════════════════════════
s = blank_slide()
bg_color(s)
page_header(s, "司机考核体系", "多维度 KPI 考核与激励约束机制")

tbl = styled_table(s, MARGIN, Inches(1.5), CONTENT_W, Inches(3.9), 8, 5)
tbl.columns[0].width = Inches(2.0)
tbl.columns[1].width = Inches(3.2)
tbl.columns[2].width = Inches(1.5)
tbl.columns[3].width = Inches(1.3)
tbl.columns[4].width = Inches(3.9)

headers = ["考核维度", "考核指标", "目标值", "权重", "考核方式"]
for i, h in enumerate(headers):
    set_cell(tbl.rows[0].cells[i], h, 12, True, WHITE, PRIMARY)

data = [
    ("时效管理", "准时到达率 / 准时交付率", "≥98%", "25%", "GPS轨迹 + 签收时间比对"),
    ("安全管理", "安全事故次数 / 违章次数", "0次", "20%", "事故记录 + 交通违章查询"),
    ("货物完好", "货物破损率 / 丢失率", "≤0.1%", "20%", "交接清点记录 + 客户反馈"),
    ("服务规范", "着装 / 话术 / 操作流程合规", "100%", "15%", "现场检查 + 客户评价"),
    ("装载规范", "缠绕膜固定 / 堆叠标准执行", "100%", "10%", "装车照片审核 + 抽查"),
    ("客户满意", "客户投诉次数 / 好评率", "0投诉", "10%", "客户回访 + 投诉记录"),
    ("油耗管理", "百公里油耗 / 路线执行率", "达标", "—", "OBD数据 + GPS轨迹分析"),
]
for r, rd in enumerate(data):
    fill = WHITE if r % 2 == 0 else LIGHT_GRAY
    for c, v in enumerate(rd):
        set_cell(tbl.rows[r + 1].cells[c], v, 10, (c == 0), fill=fill)

# Bottom: 激励 & 约束
half_w2 = (CONTENT_W - Inches(0.4)) / 2
rw = card(s, MARGIN, Inches(5.7), half_w2, Inches(1.2), fill_color=ACCENT_LIGHT, border_color=ACCENT)
tf = rw.text_frame
tf.word_wrap = True
txt(tf, "激励机制", 14, True, ACCENT)
para(tf, "月度安全之星 / 服务之星评选，给予现金奖励", 11)
para(tf, "年度优秀司机额外奖金 + 优先续约资格", 11)

pn = card(s, MARGIN + half_w2 + Inches(0.4), Inches(5.7), half_w2, Inches(1.2), fill_color=RGBColor(0xFD, 0xED, 0xEC), border_color=RED)
tf = pn.text_frame
tf.word_wrap = True
txt(tf, "约束机制", 14, True, RED)
para(tf, "连续两月考核不达标，暂停派车 / 终止合作", 11)
para(tf, "重大安全事故一票否决，立即解约并追责", 11)

page_footer(s, 7)


# ══════════════════════════════════════════════
#  S8: 异常处理机制
# ══════════════════════════════════════════════
s = blank_slide()
bg_color(s)
page_header(s, "异常处理机制", "分层分级响应，确保问题快速闭环")

levels = [
    ("一级异常（轻微）", "司机自行处理", RGBColor(0x27, 0xAE, 0x60), [
        "轻微交通拥堵（<30min）",
        "装卸等待时间略长",
        "天气影响（小雨/轻雾）",
        "",
        "司机即时上报群内，自行调整",
        "预计延误<1小时，无需升级",
    ]),
    ("二级异常（中等）", "调度 / 客服介入", ORANGE, [
        "车辆故障（可现场修复）",
        "严重拥堵（>30min）",
        "工厂装货延迟>1小时",
        "恶劣天气（暴雨/大雾）",
        "",
        "调度协调备用车辆/调整路线",
        "客服通知客户预计延误时间",
    ]),
    ("三级异常（重大）", "管理层介入", RED, [
        "交通事故（人伤/车损）",
        "货物大面积破损/丢失",
        "车辆严重故障无法行驶",
        "自然灾害/道路封闭",
        "",
        "启动应急预案，调配备用车辆",
        "管理层30min内到场/远程指挥",
        "2h内向客户提交书面报告",
    ]),
]

lv_w = Inches(3.8)
lv_gap = Inches(0.35)
for i, (level, handler, color, items) in enumerate(levels):
    x = MARGIN + (lv_w + lv_gap) * i
    hdr = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), lv_w, Inches(0.55), fill=color)
    tf = hdr.text_frame
    txt(tf, level, 14, True, WHITE, PP_ALIGN.CENTER)
    hdr.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Handler tag
    tag = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.6), Inches(2.12), lv_w - Inches(1.2), Inches(0.38),
                fill=WHITE, line=color, lw=0.75)
    tf = tag.text_frame
    txt(tf, f"处理人：{handler}", 10, False, color, PP_ALIGN.CENTER)
    tag.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Content
    body = card(s, x, Inches(2.6), lv_w, Inches(2.6), fill_color=WHITE, border_color=color)
    tf = body.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    for j, item in enumerate(items):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(10)
        p.font.color.rgb = BLACK
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(4)

# Response time bar
resp = card(s, MARGIN, Inches(5.5), CONTENT_W, Inches(1.5), fill_color=LIGHT_GRAY)
tf = resp.text_frame
tf.word_wrap = True
txt(tf, "响应时效承诺", 14, True, DARK_BG)
para(tf, "一级异常：司机即时响应，群内报备          二级异常：调度10分钟内介入，30分钟内给出方案          三级异常：管理层15分钟内介入，1小时内到场", 11)
para(tf, "所有异常处理完成后24小时内提交报告，48小时内完成闭环", 11, True)

page_footer(s, 8)


# ══════════════════════════════════════════════
#  S9: 客服保障方案
# ══════════════════════════════════════════════
s = blank_slide()
bg_color(s)
page_header(s, "客服保障方案", "无TMS系统下的线下手工运营与报表管理")

# Top banner: 运营现状
banner = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(1.5), CONTENT_W, Inches(0.75),
               fill=RGBColor(0xFF, 0xF8, 0xE1))
tf = banner.text_frame
tf.word_wrap = True
txt(tf, "运营现状：本项目无TMS系统支持，均为线下手工操作，依靠报表管理和微信群实时沟通", 13, False, ORANGE, PP_ALIGN.CENTER)
banner.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

half_w3 = Emu(int((CONTENT_W - Inches(0.4)) / 2))
right_x = Emu(int(MARGIN + half_w3 + Inches(0.4)))
section_y = Inches(2.5)

# ── Left: 沟通保障机制 ──
# Section header
tbox(s, MARGIN, section_y, half_w3, Inches(0.4),
     "沟通保障机制", 16, True, PRIMARY)

items_cm = [
    ("专属微信群", "百雀羚方 + 中通调度 + 司机 + 客服，四方同群"),
    ("实时报备", "装货 → 发车 → 在途 → 到达 → 交接，每步群内报备"),
    ("拍照留档", "装货照片、运输照片、交接照片实时群内上传"),
    ("专人对接", "配备专属客服，7×24小时在线响应"),
]
for i, (title, desc) in enumerate(items_cm):
    y = section_y + Inches(0.55) + Inches(1.05) * i
    # Number circle
    circ = shape(s, MSO_SHAPE.OVAL, MARGIN, y + Inches(0.08), Inches(0.36), Inches(0.36), fill=PRIMARY)
    tf_c = circ.text_frame
    txt(tf_c, str(i + 1), 13, True, WHITE, PP_ALIGN.CENTER)
    tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Title + desc
    tbox(s, MARGIN + Inches(0.5), y, half_w3 - Inches(0.5), Inches(0.3),
         title, 13, True, DARK_BG)
    tbox(s, MARGIN + Inches(0.5), y + Inches(0.3), half_w3 - Inches(0.5), Inches(0.3),
         desc, 11, False, DARK_GRAY)

# ── Right: 报表管理体系 ──
tbox(s, right_x, section_y, half_w3, Inches(0.4),
     "报表管理体系", 16, True, ACCENT)

tbl2 = styled_table(s, right_x, section_y + Inches(0.55), half_w3, Inches(2.3), 6, 3)
tbl2.columns[0].width = Inches(1.3)
tbl2.columns[1].width = Inches(0.8)
tbl2.columns[2].width = Emu(int(half_w3) - int(Inches(2.1)))

h2 = ["报表名称", "频率", "内容"]
for i, h in enumerate(h2):
    set_cell(tbl2.rows[0].cells[i], h, 10, True, WHITE, ACCENT)

d2 = [
    ("运输日报", "每日", "车次/货量/准时率/异常"),
    ("周度汇总", "每周", "运输量/KPI/趋势分析"),
    ("月度对账", "每月", "明细/费用/破损赔付"),
    ("异常记录", "实时", "事件/处理/结果"),
    ("司机考核", "每月", "KPI得分/排名/奖惩"),
]
for r, rd in enumerate(d2):
    fill = WHITE if r % 2 == 0 else LIGHT_GRAY
    for c, v in enumerate(rd):
        set_cell(tbl2.rows[r + 1].cells[c], v, 9, fill=fill)

# Right bottom: 保障承诺
pc = card(s, right_x, Inches(5.5), half_w3, Inches(1.2), fill_color=ACCENT_LIGHT, border_color=ACCENT)
tf = pc.text_frame
tf.word_wrap = True
txt(tf, "保障承诺", 13, True, ACCENT)
para(tf, "所有报表数据真实、准确、可追溯", 10)
para(tf, "支持随时提供原始数据与照片记录", 10)

page_footer(s, 9)


# ══════════════════════════════════════════════
#  S10: 服务保障承诺
# ══════════════════════════════════════════════
s = blank_slide()
bg_color(s, DARK_BG)

# Top accent line
shape(s, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.05), fill=ACCENT)

# Title
tbox(s, Inches(1), Inches(0.6), Inches(11.3), Inches(0.7),
     "服务保障承诺", 32, True, WHITE, PP_ALIGN.CENTER)
shape(s, MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(1.35), Inches(1.8), Pt(2), fill=ACCENT)

# Promise cards — 3×2 grid
promises = [
    ("时效保障", "准时率 ≥ 98%\n延误提前通知\n应急备用车辆"),
    ("安全保障", "全程GPS监控\n货物破损率 ≤ 0.1%\n全额货物保险"),
    ("专业团队", "专属客服 7×24h\n专业司机培训上岗\n标准化操作流程"),
    ("透明管理", "每步拍照留档\n日报/周报/月报\n数据全程可追溯"),
    ("快速响应", "异常10min响应\n1h内解决方案\n24h完成闭环"),
    ("长期合作", "持续优化改进\n定期回顾会议\n灵活调整方案"),
]

card_w = Inches(3.5)
card_h = Inches(1.7)
card_gap_x = Inches(0.5)
card_gap_y = Inches(0.4)
start_x = (SLIDE_W - card_w * 3 - card_gap_x * 2) / 2
start_y = Inches(1.8)

for i, (title, desc) in enumerate(promises):
    cx = start_x + (card_w + card_gap_x) * (i % 3)
    cy = start_y + (card_h + card_gap_y) * (i // 3)
    c = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, card_h,
              fill=RGBColor(0x14, 0x3D, 0x5C))
    tf = c.text_frame
    tf.word_wrap = True
    txt(tf, title, 16, True, WHITE, PP_ALIGN.CENTER)
    for line in desc.split("\n"):
        para(tf, line, 11, False, RGBColor(0xA8, 0xC8, 0xE8), PP_ALIGN.CENTER, sb=Pt(3))

# Bottom: tagline + logos (proper spacing)
tagline_y = Inches(6.0)
tbox(s, Inches(1), tagline_y, Inches(11.3), Inches(0.5),
     "期待与百雀羚携手共创高效物流体系", 18, True, WHITE, PP_ALIGN.CENTER)

# Logo row — properly sized and centered
logo_row_y = Inches(6.55)
# ZTO logo: small, fit within bounds
add_img(s, ZTO_LOGO, Inches(5.3), logo_row_y, w=Inches(0.55))
# Company name next to logo
tbox(s, Inches(6.0), logo_row_y + Inches(0.08), Inches(3.5), Inches(0.35),
     "中通快递股份有限公司", 13, False, RGBColor(0xA8, 0xC8, 0xE8))

# Bottom bar
shape(s, MSO_SHAPE.RECTANGLE, 0, Inches(7.3), SLIDE_W, Inches(0.2), fill=ACCENT)


# ══════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")

# Auto-upload to Feishu (overwrite existing)
import subprocess, tempfile, shutil
FEISHU_FOLDER = "AgH1fKkDOleiKSdQWJ1cN67fnuc"
FILENAME = os.path.basename(OUTPUT)
print(f"\nUploading to Feishu...")
# Stage only the pptx in a temp subdir so +push doesn't upload other files
staging = os.path.join(SCRIPT_DIR, ".upload_staging")
os.makedirs(staging, exist_ok=True)
shutil.copy2(OUTPUT, os.path.join(staging, FILENAME))
r = subprocess.run(
    ["lark-cli", "drive", "+push",
     "--folder-token", FEISHU_FOLDER,
     "--local-dir", ".upload_staging",
     "--if-exists", "overwrite"],
    cwd=SCRIPT_DIR, capture_output=True, text=True
)
shutil.rmtree(staging, ignore_errors=True)
if r.returncode == 0:
    import json
    result = json.loads(r.stdout)
    if result.get("ok"):
        items = result["data"].get("items", [])
        for it in items:
            if it["rel_path"] == FILENAME:
                print(f"Uploaded: https://fqmtvue07d8.feishu.cn/file/{it['file_token']}")
                break
        else:
            print(f"Uploaded (token not matched)")
    else:
        print(f"Upload failed: {result.get('error', {}).get('message', 'unknown')}")
else:
    print(f"Upload error: {r.stderr.strip()}")
