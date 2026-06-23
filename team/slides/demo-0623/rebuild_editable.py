# -*- coding: utf-8 -*-
"""
将【打破工程困局_仿真分享_30min.pptx】中的 15 张贴图幻灯片
转换为可编辑的原生 PPT 内容（文本框 + 形状 + 表格）。

- 贴图页（PPTX 第 1,3,4,5,6,7,8,9,10,11,12,13,14,19,21 页）来源于 HTML deck 的 15 页，
  逐页移除图片并用原生形状重建。
- 原生页（第 2,15,16,17,18,20 页）保持不变。
- 配色 / 字体严格对齐已有原生页，保证整套视觉风格一致。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SRC = "打破工程困局_仿真分享_30min.原始贴图版.pptx"  # 始终从原始贴图版重建，保证幂等
OUT = "打破工程困局_仿真分享_30min.pptx"

# ---------------- 调色板（对齐原生页 srgbClr）----------------
BG      = RGBColor(0x0A, 0x0E, 0x18)
PANEL   = RGBColor(0x12, 0x1A, 0x2B)
PANEL2  = RGBColor(0x0E, 0x14, 0x1F)
BORDER  = RGBColor(0x2A, 0x3A, 0x55)
TEXT    = RGBColor(0xEC, 0xF1, 0xF9)
MUTED   = RGBColor(0x97, 0xA3, 0xB8)
DIM     = RGBColor(0x5A, 0x66, 0x78)
ORANGE  = RGBColor(0xFF, 0x8C, 0x28)
CYAN    = RGBColor(0x4F, 0xB6, 0xFF)
GREEN   = RGBColor(0x3F, 0xD6, 0x8C)
RED     = RGBColor(0xFF, 0x6B, 0x6B)
LEGACY  = RGBColor(0x8A, 0x97, 0xA8)

MONO = "Consolas"
CN   = "Microsoft YaHei"

L_MARGIN = 0.5
USABLE_L = 0.72
USABLE_R = 12.83
USABLE_W = USABLE_R - USABLE_L  # 12.11


# ---------------- 基础形状 / 文本 helper ----------------
def add_box(slide, l, t, w, h, fill=None, line=None, line_w=1.0, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    if shp.has_text_frame:
        shp.text_frame.word_wrap = True
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(shp.text_frame, m, Emu(0))
    return shp


def add_text(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Emu(0))
    tf.vertical_anchor = anchor
    return tb, tf


def para(tf, runs, size=14, color=TEXT, bold=False, font=CN, align=PP_ALIGN.LEFT,
         first=False, space_after=0, space_before=0, line=1.0):
    """runs: str 或 [(text, {opts})...]，opts 可含 size/bold/color/font。"""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if space_after:
        p.space_after = Pt(space_after)
    if space_before:
        p.space_before = Pt(space_before)
    try:
        p.line_spacing = line
    except Exception:
        pass
    if isinstance(runs, str):
        runs = [(runs, {})]
    for txt, opts in runs:
        r = p.add_run(); r.text = txt
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.name = opts.get("font", font)
        r.font.color.rgb = opts.get("color", color)
    return p


# ---------------- HUD 框架组件 ----------------
def hud_top(slide, left, right=""):
    _, tf = add_text(slide, L_MARGIN, 0.2, 8.0, 0.3)
    para(tf, [(left, {})], size=10, color=MUTED, font=MONO, first=True)
    if right:
        _, tf2 = add_text(slide, 4.9, 0.2, 7.93, 0.3)
        para(tf2, [(right, {})], size=10, color=CYAN, font=MONO,
             align=PP_ALIGN.RIGHT, first=True)


def hud_bottom(slide, left, right=""):
    add_box(slide, L_MARGIN, 7.05, 12.33, 0.012, fill=BORDER)
    _, tf = add_text(slide, L_MARGIN, 7.12, 7.5, 0.3)
    para(tf, [(left, {})], size=9, color=DIM, font=MONO, first=True)
    if right:
        _, tf2 = add_text(slide, 4.9, 7.12, 7.93, 0.3)
        para(tf2, [(right, {})], size=9, color=DIM, font=MONO,
             align=PP_ALIGN.RIGHT, first=True)


def title(slide, runs, t=0.55, size=28):
    """runs: [(text,{color/bold})...]；左侧加橙色竖条。"""
    add_box(slide, L_MARGIN, t + 0.04, 0.09, 0.52, fill=ORANGE)
    _, tf = add_text(slide, USABLE_L, t, 12.0, 0.7)
    para(tf, runs, size=size, color=TEXT, bold=True, first=True, line=1.05)


def lead(slide, runs, t=1.32, w=12.0, color=MUTED, size=14.5):
    _, tf = add_text(slide, USABLE_L, t, w, 0.7)
    para(tf, runs, size=size, color=color, first=True, line=1.25)


def panel(slide, l, t, w, h, accent=None, fill=PANEL, border=BORDER):
    sp = add_box(slide, l, t, w, h, fill=fill, line=border, line_w=1)
    if accent:
        add_box(slide, l, t, w, 0.055, fill=accent)
    return sp


def set_bg(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG


def clear_pictures(slide):
    n = 0
    for shp in list(slide.shapes):
        if shp.shape_type == 13:  # PICTURE
            shp._element.getparent().remove(shp._element)
            n += 1
    # 移除孤立的图片关系，使图片不再被打包进文件
    part = slide.part
    drop = [rid for rid, rel in list(part.rels.items())
            if rel.reltype.endswith("/image")]
    for rid in drop:
        try:
            part.drop_rel(rid)
        except Exception:
            pass
    return n


# ---------------- 复合布局组件 ----------------
def cols(n, gap=0.3, l=USABLE_L, r=USABLE_R):
    total = r - l
    w = (total - gap * (n - 1)) / n
    return [(l + i * (w + gap), w) for i in range(n)]


def terminal(slide, l, t, w, h, bar, lines):
    """代码终端块。lines: [run-list, ...]"""
    panel(slide, l, t, w, h, fill=PANEL2, border=BORDER)
    add_box(slide, l, t, w, 0.3, fill=RGBColor(0x18, 0x20, 0x30))
    for i, c in enumerate([RED, ORANGE, GREEN]):
        add_box(slide, l + 0.14 + i * 0.16, t + 0.11, 0.08, 0.08, fill=c, rounded=True)
    _, tb = add_text(slide, l + 0.62, t + 0.03, w - 0.72, 0.24)
    para(tb, [(bar, {})], size=9.5, color=MUTED, font=MONO, first=True)
    _, tf = add_text(slide, l + 0.22, t + 0.42, w - 0.44, h - 0.5)
    for i, segs in enumerate(lines):
        para(tf, segs, size=10.5, font=MONO, color=MUTED, first=(i == 0), line=1.32)


def card_block(slide, l, t, w, h, head, body, accent=CYAN, head_color=TEXT,
               head_size=13.5, body_size=11.5, en=None):
    panel(slide, l, t, w, h, fill=PANEL)
    add_box(slide, l, t, 0.05, h, fill=accent)
    _, tf = add_text(slide, l + 0.22, t + 0.16, w - 0.4, h - 0.3)
    hruns = [(head, {})]
    if en:
        hruns.append(("  " + en, {"color": DIM, "size": head_size - 3, "bold": False}))
    para(tf, hruns, size=head_size, color=head_color, bold=True, first=True, space_after=5)
    if body:
        para(tf, body if isinstance(body, list) else [(body, {})],
             size=body_size, color=MUTED, line=1.22)


def mini_card(slide, l, t, w, h, head, body, accent=CYAN):
    panel(slide, l, t, w, h, accent=accent)
    _, tf = add_text(slide, l + 0.2, t + 0.22, w - 0.4, h - 0.38)
    para(tf, [(head, {})], size=13, color=TEXT, bold=True, first=True, space_after=5)
    para(tf, [(body, {})], size=11, color=MUTED, line=1.22)


def cmp_panel(slide, l, t, w, h, accent, head, head_color, rows, val_color):
    panel(slide, l, t, w, h, accent=accent)
    _, tf = add_text(slide, l + 0.25, t + 0.24, w - 0.5, h - 0.4)
    para(tf, [(head, {})], size=15, color=head_color, bold=True, first=True, space_after=10)
    for k, v in rows:
        para(tf, [(k + "   ", {"color": DIM, "size": 11}),
                  (v, {"color": val_color, "size": 12.5})],
             space_after=7, line=1.1)


def simple_table(slide, l, t, w, rows_data, col_w, header_fill, row_h=0.55,
                 head_h=0.5, fs=11.5, head_fs=12, cell_colors=None):
    """rows_data: list of rows; row[0] 为表头。col_w: 比例列表。"""
    nrows = len(rows_data)
    ncols = len(rows_data[0])
    total_h = head_h + row_h * (nrows - 1)
    gtbl = slide.shapes.add_table(nrows, ncols, Inches(l), Inches(t),
                                  Inches(w), Inches(head_h + row_h * (nrows - 1)))
    tbl = gtbl.table
    tbl.first_row = False
    tbl.horz_banding = False
    cw = sum(col_w)
    for j, ratio in enumerate(col_w):
        tbl.columns[j].width = Inches(w * ratio / cw)
    tbl.rows[0].height = Inches(head_h)
    for i in range(1, nrows):
        tbl.rows[i].height = Inches(row_h)
    for i, row in enumerate(rows_data):
        for j, cell_text in enumerate(row):
            cell = tbl.cell(i, j)
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                col, bold, size = TEXT, True, head_fs
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PANEL if i % 2 else PANEL2
                col, bold, size = MUTED, False, fs
                if cell_colors:
                    col = cell_colors(i, j)
            tfc = cell.text_frame
            tfc.word_wrap = True
            p = tfc.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = cell_text
            r.font.size = Pt(size); r.font.bold = bold
            r.font.name = CN; r.font.color.rgb = col
    return tbl


# ====================================================================
#  15 页逐页重建（对应 HTML deck 1..15）
# ====================================================================
def s01_title(s):
    hud_top(s, "# PROJECT_LEAD: 仿真算法负责人",
            "# TARGET: 2026 Intern Technical Onboarding")
    _, tf = add_text(s, 0.8, 2.0, 11.73, 1.2, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("打破工程困局", {})], size=58, color=TEXT, bold=True,
         align=PP_ALIGN.CENTER, first=True)
    _, tf2 = add_text(s, 1.5, 3.35, 10.33, 1.0)
    para(tf2, [("如何在", {}), ("“不确定”", {"color": ORANGE, "bold": True}),
               ("的物理世界里，为自动驾驶大模型构建", {}),
               ("“确定”", {"color": ORANGE, "bold": True}), ("的运行闭环", {})],
         size=18, color=MUTED, align=PP_ALIGN.CENTER, first=True, line=1.35)
    add_box(s, 5.05, 4.75, 0.12, 0.12, fill=CYAN, rounded=True)
    _, tf3 = add_text(s, 5.3, 4.62, 6.0, 0.4)
    para(tf3, [("仿真算法负责人：", {"color": MUTED}),
               ("Li Kun", {"color": TEXT, "bold": True})],
         size=15, first=True)
    _, tf4 = add_text(s, 0.8, 5.45, 11.73, 0.4)
    para(tf4, [("XPENG AUTOMOTIVE SIMULATION ALGORITHM TEAM", {})],
         size=11, color=DIM, font=MONO, align=PP_ALIGN.CENTER, first=True)
    hud_bottom(s, "RFC-SIM-2024-V2.0 // PROPOSAL / INTERNAL ONLY",
               "SYSTEM_STATUS: ONLINE")


def s02_paradigm(s):
    hud_top(s, "", "DATA.STREAM: VLA-2.0-LIVE")
    title(s, [("范式转移：VLA 2.0 端到端架构拿掉了过去的", {}),
              ("“安全拐杖”", {"color": ORANGE})])
    lead(s, [("端到端架构之后，我们没法再靠看代码（静态 Review）来保证安全性了。", {})])
    (l1, w1), (l2, w2) = cols(2)
    cmp_panel(s, l1, 1.95, w1, 2.15, LEGACY, "传统模块化（Legacy）", MUTED,
              [("核心输入", "BBox / 车道线轨迹"),
               ("逻辑载体", "C++ Hard Rules"),
               ("演进路径", "人工修 Bug")], LEGACY)
    cmp_panel(s, l2, 1.95, w2, 2.15, CYAN, "VLA 2.0（端到端）", CYAN,
              [("核心输入", "原始视频流 / 点云流"),
               ("逻辑载体", "深度神经网络参数"),
               ("演进路径", "数据闭环驱动参数迭代")], TEXT)
    for (cl, cw), (h, b) in zip(cols(3), [
        ("看不见的黑盒", "整套驾驶逻辑都压进了神经网络权重里，没有 BBox、没有规则，既看不见也改不动。"),
        ("没有插桩点", "感知到决策一步到位（<80ms），中间不再有接口，让我们断言“它到底有没有看到行人”。"),
        ("行为不可预判", "能在无标线路段类人避让，也可能突然失灵——光读代码，无法预测它的能力边界。")]):
        mini_card(s, cl, 4.35, cw, 2.0, h, b)
    hud_bottom(s, "SYS.STATUS: ACTIVE // DATA.STREAM: VLA-2.0-LIVE", "LATENCY: <80MS")


def s03_dataloop(s):
    hud_top(s, "SYS.STATUS: ACTIVE // LATENCY: <80MS",
            "DATA.STREAM: VLA-2.0-LIVE // DATE: 2024-05-15")
    title(s, [("数据闭环：VLA 模型进化的唯一", {}), ("“强核燃料”", {"color": ORANGE})])
    _, tf = add_text(s, USABLE_L, 1.6, 5.9, 1.2)
    para(tf, [("仿真不仅是测试工具，更是车端 VLA 2.0 模型进行", {}),
              ("强化学习（RL）", {"color": CYAN, "bold": True}),
              ("的虚拟教练。", {})], size=16, color=TEXT, first=True, line=1.3)
    terminal(s, USABLE_L, 3.0, 5.9, 2.3, "scaling_law.py",
             [[("# Scaling Law of Simulation", {"color": DIM})],
              [("Model_Intelligence =", {}), (" Data_Diversity", {"color": CYAN})],
              [("  * Loop_Iterations * Sim_Realism", {"color": CYAN})],
              [("Target_Cost = ", {}), ("0.01", {"color": ORANGE}),
               (" * Real_World_Cost", {})]])
    rx = 6.92
    _, tt = add_text(s, rx, 1.6, 5.9, 0.4)
    para(tt, [("Flywheel Loop", {})], size=13, color=ORANGE, bold=True,
         font=MONO, first=True)
    fw = [("大规模并发（Massive Mining）", "每日千万级里程并行验证。仿真系统是版本准出的 Gatekeeper。"),
          ("Case 集验证（Variant Bombing）", "自动提取路测 Hard Case，泛化 100 倍变体场景进行极限轰炸。"),
          ("版本筛选（Filtering）", "筛选正向收益模型权重，剔除端到端黑盒导致的逻辑退化。")]
    for i, (h, b) in enumerate(fw):
        card_block(s, rx, 2.05 + i * 1.12, 5.9, 1.0, h, b, accent=CYAN,
                   head_size=12.5, body_size=10.5)
    hud_bottom(s, "SYS.STATUS: ACTIVE // DATA.STREAM: VLA-2.0-LIVE",
               "LATENCY: <80MS // DATE: 2024-05-15")


def s04_deadlock(s):
    hud_top(s, "DATA.STREAM: VLA-LEGACY-SIM-CRISIS // SYS.STATUS: ACTIVE",
            "LATENCY: <80NS")
    title(s, [("技术债危机：传统仿真架构面对 VLA 的", {}), ("“双重死局”", {"color": ORANGE})])
    p = panel(s, USABLE_L, 1.55, USABLE_W, 0.62, fill=PANEL2)
    _, tf = add_text(s, USABLE_L + 0.25, 1.66, USABLE_W - 0.5, 0.45,
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("面对 VLA 端到端，我们必须找到一种", {}),
              ("“既有 LogSim 真实度，又有 WorldSim 交互性”", {"color": CYAN, "bold": True}),
              ("的新方案。", {})], size=14, color=TEXT, first=True,
         align=PP_ALIGN.CENTER)
    cards = [("LogSim（数据回放）", "回放实车录制的传感器 Log。",
              "开环（Open-loop）。环境不随本车动作改变。VLA 内部没有 BBox 接口，我们根本无法 Assert“是否检测到行人”。"),
             ("WorldSim（白模规则）", "手工搭建 3D 资产，编写 NPC 规则。",
              "Sim2Real Gap。画面因为过于“整洁”而严重失真，视觉大模型根本不认 WorldSim 生成的“游戏画面”。")]
    for (cl, cw), (h, mech, weak) in zip(cols(2), cards):
        panel(s, cl, 2.4, cw, 4.0, accent=RED)
        _, tf = add_text(s, cl + 0.28, 2.62, cw - 0.56, 3.6)
        para(tf, [(h, {})], size=18, color=TEXT, bold=True, first=True, space_after=12)
        para(tf, [("机制   ", {"color": CYAN, "bold": True, "size": 11}),
                  (mech, {"size": 13})], color=TEXT, space_after=12, line=1.3)
        para(tf, [("致命弱点   ", {"color": ORANGE, "bold": True, "size": 11}),
                  (weak, {"size": 13})], color=MUTED, line=1.4)
        bt = add_box(s, cl + 0.28, 5.85, cw - 0.56, 0.42, fill=RGBColor(0x2A, 0x14, 0x14))
        _, btf = add_text(s, cl + 0.28, 5.9, cw - 0.56, 0.32, anchor=MSO_ANCHOR.MIDDLE)
        para(btf, [("DEAD END: VLA INCOMPATIBLE", {})], size=11, color=RED,
             bold=True, font=MONO, align=PP_ALIGN.CENTER, first=True)
    hud_bottom(s, "SYS.STATUS: ACTIVE", "LATENCY: <80MS // DATE: 2024-05-15")


def s05_spatiotemporal(s):
    hud_top(s, "DATA.STREAM: SPATIOTEMPORAL-EDITING // SYS.STATUS: ACTIVE // LATENCY: <10MS",
            "SYSTEM_MODE = REALTIME_INTERVENTION")
    title(s, [("时空连续性编辑：从", {}), ("“录像重放”", {"color": ORANGE}),
              ("到", {}), ("“无限变体考场”", {"color": ORANGE})])
    lead(s, [("借助自研 3DGS 引擎结合时空连续性算法，我们能对场景中的每一个 Gaussian 实体进行实时干预，"
              "彻底解决纯 2D 扩散模型视角闪烁与畸变的痛点。", {})], w=12.0)
    panel(s, USABLE_L, 2.35, 2.5, 3.9, fill=PANEL2, border=CYAN)
    _, tf = add_text(s, USABLE_L, 2.35, 2.5, 3.9, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("⎇", {})], size=40, color=CYAN, align=PP_ALIGN.CENTER, first=True,
         space_after=6)
    para(tf, [("Baseline", {})], size=13, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    para(tf, [("Log Root", {})], size=13, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    items = [("NPC 变轨", "(Trajectory Hack)",
              "抽离原始 Log 里的正常行驶车辆，修改轨迹，令其对本车执行“强行 Cut-in”（闭环博弈一致性）。"),
             ("资产替换", "(Asset Swap)",
              "像素级抠除路边的一棵树，无缝替换为横穿马路的行人。"),
             ("无限繁衍", "(Infinite Spawning)",
              "在同一地段瞬间生成上百个差异化的高保真 Corner Case。")]
    bx = 3.6
    for i, (h, en, b) in enumerate(items):
        card_block(s, bx, 2.35 + i * 1.33, USABLE_R - bx, 1.18, h, b,
                   accent=ORANGE, en=en, head_size=14, body_size=12)
    hud_bottom(s, "LATENCY: <10MS", "DATE: 2024-05-15")


def s06_covariate(s):
    hud_top(s, "DATA.STREAM: COVARIATE-SHIFT-CASCADE // SYS.STATUS: ACTIVE // LATENCY: <80MS", "")
    title(s, [("深度推演：协变量偏移（Covariate Shift）灾难链条", {})])
    steps = [("T = 0s", "本车 Action 与原始记录产生微小的 0.1° 偏差。", CYAN),
             ("T = 5s", "偏差由于时间累积，车辆实际物理位置已偏离 1.5 米！", ORANGE),
             ("T = 5s+ 空间错乱", "视野与录像中的背景发生严重空间错乱（Spatial Misalignment）。", ORANGE),
             ("Result · 模型崩溃", "看到的画面与其所处的动作逻辑完全脱节，无法测试 VLA 偏差后的纠偏能力。", RED)]
    cs = cols(4, gap=0.28)
    for (cl, cw), (t, b, c) in zip(cs, steps):
        panel(s, cl, 1.75, cw, 2.7, accent=c)
        _, tf = add_text(s, cl + 0.2, 1.98, cw - 0.4, 2.4)
        para(tf, [(t, {})], size=14, color=c, bold=True, font=MONO, first=True,
             space_after=10)
        para(tf, [(b, {})], size=12, color=MUTED, line=1.3)
    for i in range(3):
        _, ar = add_text(s, cs[i][0] + cs[i][1] - 0.02, 2.9, 0.3, 0.4)
        para(ar, [("▶", {})], size=14, color=DIM, align=PP_ALIGN.CENTER, first=True)
    panel(s, USABLE_L, 4.75, USABLE_W, 1.7, fill=PANEL2)
    _, tf = add_text(s, USABLE_L + 0.3, 4.98, USABLE_W - 0.6, 1.3)
    para(tf, [("黑盒测试的唯一出口：", {"color": ORANGE, "bold": True}),
              ("我们只能通过", {}), ("闭环（Closed-loop）", {"color": ORANGE, "bold": True}),
              ("来观察——如果不接管，车会不会撞？", {})],
         size=14, color=TEXT, first=True, space_after=10, line=1.3)
    para(tf, [("仿真底线：", {"color": ORANGE, "bold": True}),
              ("必须提供", {}), ("实时渲染", {"color": ORANGE, "bold": True}),
              ("，以", {}), ("毫秒级响应", {"color": ORANGE, "bold": True}),
              (" VLA 的每一个转角输出。", {})], size=14, color=TEXT, line=1.3)
    hud_bottom(s, "SYS.STATUS: ACTIVE / DATA.STREAM: COVARIATE-SHIFT-CASCADE",
               "LATENCY: <80MS / DATE: 2024-05-15")


def s07_nerf(s):
    hud_top(s, "DATA.STREAM: NERF-3DGS-EVOLUTION // SYS.STATUS: ACTIVE // LATENCY: <10MS", "")
    title(s, [("神经场景基石：从隐式 NeRF 到显式 3DGS 的质变", {})])
    rows = [["技术", "表示形式", "存储成本", "渲染延迟", "编辑性"],
            ["NeRF", "MLP Weights（隐式）", "~50MB（极小）", "1-2s / 帧（极慢）",
             "极难（Implicit）→ 无法支持 20Hz 闭环"],
            ["3DGS", "Explicit Splats（显式）", "~300MB（较大）", "100+ FPS（极速）",
             "易编辑（Point-based）→ 照片级实时交互"]]

    def cc(i, j):
        if j == 0:
            return RED if i == 1 else GREEN
        return LEGACY if i == 1 else TEXT
    simple_table(s, USABLE_L, 1.7, USABLE_W, rows, [0.12, 0.22, 0.18, 0.2, 0.32],
                 PANEL2, row_h=0.95, head_h=0.5, fs=12, cell_colors=cc)
    terminal(s, 6.6, 4.5, 6.23, 1.95, "# Gaussian Map Profiling",
             [[("Total_Splat_Count = ", {}), ("2,500,000", {"color": ORANGE})],
              [("Memory_Usage = ", {}), ("420 MB", {"color": ORANGE})],
              [("Inference_Time_Per_View = ", {}), ("8.5ms", {"color": ORANGE})],
              [("System_Status = ", {}), ("READY_FOR_CLOSED_LOOP", {"color": GREEN})]])
    _, tf = add_text(s, USABLE_L, 4.7, 5.7, 1.6)
    para(tf, [("结论", {})], size=13, color=CYAN, bold=True, font=MONO, first=True,
         space_after=8)
    para(tf, [("3DGS 以可接受的存储代价，换来两个数量级的渲染提速与可编辑性——"
               "这是支撑实时闭环仿真的关键基石。", {})],
         size=13, color=MUTED, line=1.35)
    hud_bottom(s, "SYS.STATUS: ACTIVE", "LATENCY: <10MS // DATE: 2024-05-15")


def s08_pipeline(s):
    hud_top(s, "DATA.STREAM: DRIVING-GAUSSIAN-PIPELINE // SYS.STATUS: AUTOMATED // LATENCY: <50MS", "")
    title(s, [("工业化管线：DrivingGaussian 大规模自动化重建", {})])
    steps = [("数据对齐", "(Data Alignment)", "Camera + LiDAR 高精度时空同步。", CYAN, False),
             ("静态背景构建", "(Static Core)", "LiDAR Prior 引导，利用点云滤除视觉幽灵噪点。", CYAN, False),
             ("动态物体绑定", "(Dynamic Binding)", "核心难点突破！将 Gaussian 点绑定至运动轨迹，实现完美“动静分离”。", ORANGE, True),
             ("实时光栅化", "(Rasterization)", "1080P / 120FPS 极速出图。", CYAN, False)]
    cw = (USABLE_W - 3 * 0.5) / 4
    for i, (h, en, b, c, key) in enumerate(steps):
        cl = USABLE_L + i * (cw + 0.5)
        panel(s, cl, 2.0, cw, 3.0, accent=c, fill=(RGBColor(0x1A, 0x16, 0x10) if key else PANEL))
        _, tf = add_text(s, cl + 0.18, 2.24, cw - 0.36, 2.6)
        para(tf, [(h, {})], size=14, color=(ORANGE if key else TEXT), bold=True,
             first=True)
        para(tf, [(en, {})], size=9.5, color=DIM, font=MONO, space_after=10)
        para(tf, [(b, {})], size=11.5, color=MUTED, line=1.3)
        if i < 3:
            _, ar = add_text(s, cl + cw + 0.06, 3.2, 0.38, 0.5)
            para(ar, [("▶", {})], size=16, color=ORANGE, align=PP_ALIGN.CENTER, first=True)
    panel(s, USABLE_L, 5.35, USABLE_W, 1.15, fill=PANEL2, border=ORANGE)
    _, tf = add_text(s, USABLE_L + 0.3, 5.55, USABLE_W - 0.6, 0.8,
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("ROI：", {"color": ORANGE, "bold": True}),
              ("将单个街区的重建时间从 ", {}), ("2 周", {"color": ORANGE, "bold": True, "size": 16}),
              ("（依赖美术人工）瞬间压缩至 ", {}), ("2 小时", {"color": ORANGE, "bold": True, "size": 16}),
              ("（全自动流水线）。场景还原度 PSNR 突破 ", {}),
              ("32dB+", {"color": ORANGE, "bold": True, "size": 16}), ("。", {})],
         size=13.5, color=TEXT, first=True, line=1.35)
    hud_bottom(s, "LATENCY: <50MS / DATE: 2024-05-15",
               "SYSTEM_MODE = FULL_AUTO_RECONSTRUCTION")


def s09_gensim(s):
    hud_top(s, "DATA.STREAM: WORLD-MODEL-PREDICTION // SYS.STATUS: ACTIVE // LATENCY: <15MS", "")
    title(s, [("Gen-Sim 世界模型：为环境注入预测未来的", {}), ("“脑力”", {"color": ORANGE})])
    lead(s, [("环境不再是被动的静态贴图，而是具备物理常识的智能体集合。"
              "这是解决“仿真太假”与“NPC 太呆”的终极武器。", {})], w=12.0)
    panel(s, 5.16, 2.5, 3.0, 1.9, fill=PANEL2, border=CYAN)
    _, tf = add_text(s, 5.16, 2.5, 3.0, 1.9, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("🌐", {})], size=34, color=CYAN, align=PP_ALIGN.CENTER, first=True,
         space_after=4)
    para(tf, [("Digital World", {})], size=12, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    para(tf, [("Model Node", {})], size=12, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    cards = [("时空预测", "(Spatiotemporal)", "给定 {当前帧 + 本车 Action}，精准推演未来 T 秒的连续场景变化。", USABLE_L, 2.4),
             ("因果推理", "(Causality)", "打破死板逻辑树。“如果我急刹车，后方 NPC 必定按物理常识做出减速博弈反应。”", 9.0, 2.4),
             ("万物多样性", "(Diversity)", "对同一路口，自由且符合常理地生成雨、雾、黄昏等多重光影与气候状态。", USABLE_L, 4.7)]
    for h, en, b, cl, ct in cards:
        card_block(s, cl, ct, 3.83, 2.0, h, b, accent=CYAN, en=en,
                   head_size=14, body_size=11.5)
    card_block(s, 9.0, 4.7, 3.83, 2.0,
               "RL 训练土壤", "高保真、可交互、可泛化——世界模型为 VLA 的强化学习提供唯一可规模化的训练环境。",
               accent=ORANGE, head_size=14, body_size=11.5)
    hud_bottom(s, "LATENCY: <15MS // DATE: 2024-05-15",
               "SYSTEM_MODE = PREDICTIVE_GEN-SIM")


def s10_dimension(s):
    hud_top(s, "DATA.STREAM: WORLD-MODEL-COMPARISON // SYS.STATUS: ACTIVE // LATENCY: <15MS", "")
    title(s, [("维度碾压：传统白模规则 vs AI 生成式世界模型", {})])
    rows = [["核心维度", "传统规则引擎（WorldSim）", "AI 世界模型（Gen-Sim）"],
            ["场景生成", "人工白模 / 样条线铺设（耗时极长）", "文本或参考视频秒级生成（极速）"],
            ["NPC 行为", "固化逻辑树（If-Else，易死锁）", "基于物理常识的自博弈（类人）"],
            ["长尾覆盖", "工程师手动硬刚（挂一漏万）", "AI 自动泛化大量边缘场景（涌现）"],
            ["物理反馈", "粗糙刚体碰撞（缺乏传感器细节）", "精确复现运动模糊、雨天反光等细节"]]

    def cc(i, j):
        if j == 0:
            return TEXT
        return LEGACY if j == 1 else GREEN
    simple_table(s, USABLE_L, 1.65, USABLE_W, rows, [0.18, 0.41, 0.41],
                 PANEL2, row_h=0.62, head_h=0.5, fs=12.5, cell_colors=cc)
    # 列着色：自定义二次写入（simple_table 默认 muted，这里重画值列颜色）
    panel(s, USABLE_L, 5.65, USABLE_W, 0.85, fill=PANEL2, border=CYAN)
    _, tf = add_text(s, USABLE_L + 0.3, 5.74, USABLE_W - 0.6, 0.66,
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("“世界模型不是在‘播放视频’，而是在‘即时构思物理现实’——"
               "这是 VLA 模型进行 RL 训练的唯一高保真土壤。”", {})],
         size=14, color=TEXT, bold=True, first=True, align=PP_ALIGN.CENTER, line=1.3)
    hud_bottom(s, "DATE: 2024-05-15", "SYSTEM_MODE = PREDICTIVE_GEN-SIM")


def s11_topology(s):
    hud_top(s, "DATA.STREAM: HYBRID-ARCH-TOPO // SYS.STATUS: ACTIVE // LATENCY: <5MS",
            "SYSTEM_MODE = PREDICTIVE_GEN-SIM")
    title(s, [("终极拓扑：3DGS 与 DiT 生成式模型的最强混合架构", {})])
    card_block(s, USABLE_L, 1.7, 5.9, 1.45,
               "预测生成引擎（DiT Engine）", "提供“时空生成能力”与底层物理因果推理。",
               accent=ORANGE, head_size=15, body_size=12.5)
    card_block(s, USABLE_L, 3.3, 5.9, 1.45,
               "基础空间层（3DGS Base）", "提供极致“显式空间一致性”与极速实时光栅化能力。",
               accent=CYAN, head_size=15, body_size=12.5)
    terminal(s, 6.92, 1.7, 5.9, 3.05, "TERMINAL",
             [[("Input: ", {}), ("{3DGS_Scene, Ego_Action}", {"color": DIM})],
              [("Process: DiT predicts", {})],
              [("  Delta_Gaussian_Parameters", {"color": CYAN})],
              [("Output:", {})],
              [("  Dynamic Closed-loop", {"color": ORANGE})],
              [("  Render (120 FPS)", {"color": ORANGE})]])
    panel(s, USABLE_L, 5.0, USABLE_W, 1.45, fill=PANEL2, border=ORANGE)
    _, tf = add_text(s, USABLE_L + 0.3, 5.2, USABLE_W - 0.6, 1.05,
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("破局意义：", {"color": ORANGE, "bold": True}),
              ("纯 2D Sora 路线无法保证多摄像头的严格对齐。通过在 ", {}),
              ("3D 特征空间", {"color": CYAN, "bold": True}),
              ("直接进行生成，我们彻底解决了环视拼缝失真的工业级难题。", {})],
         size=14, color=TEXT, first=True, line=1.4)
    hud_bottom(s, "DATA.STREAM: HYBRID-ARCH-TOPO // SYS.STATUS: ACTIVE // LATENCY: <5MS",
               "SYSTEM_MODE = PREDICTIVE_GEN-SIM // DATE: 2024-05-16")


def s12_vlm(s):
    hud_top(s, "DATA.STREAM: VLM-ASSET-REFINERY // SYS.STATUS: ACTIVE // LATENCY: <8MS",
            "SYSTEM_MODE = AUTOMATED_LOG_PROCESSING")
    title(s, [("AI 重塑管线 01：VLM（视觉大模型）自动化资产沉淀", {})])
    lead(s, [("核心痛点：", {"color": CYAN, "bold": True}),
             ("人力根本无法 Review 每天产生的 TB 级路测 Log 录像。", {})], w=12.0)
    box = panel(s, USABLE_L, 2.05, 5.9, 0.55, fill=PANEL2, border=CYAN)
    _, tf = add_text(s, USABLE_L + 0.2, 2.05, 5.5, 0.55, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("Industrial Refinery · Raw Logs  →  VLM Scanner", {})],
         size=11.5, color=CYAN, font=MONO, first=True)
    card_block(s, USABLE_L, 2.8, 5.9, 1.5, "场景理解（Semantic Understanding）",
               "自动扫描并输出“无保护左转”“复杂路口博弈”等高价值语义标签。",
               accent=CYAN, head_size=13.5, body_size=12)
    card_block(s, USABLE_L, 4.45, 5.9, 1.5, "缺陷发现（Defect Discovery）",
               "自动寻找 VLA 表现不佳（如急刹车）的薄弱片段，直接送入飞轮供模型重新训练。",
               accent=ORANGE, head_size=13.5, body_size=12)
    terminal(s, 6.92, 2.05, 5.9, 3.9, "Agent Log",
             [[("Agent Log:", {})],
              [("> Analyzing Video Segment #1024...", {})],
              [("> Scenario Detected:", {})],
              [("  \"Aggressive Cut-in by Truck\"", {"color": ORANGE})],
              [("> VLA Performance:", {})],
              [("  \"Late Braking (TTC=0.8s)\"", {"color": ORANGE})],
              [("> Recommendation:", {})],
              [("  Extract to High-Priority Set.", {"color": GREEN})]])
    hud_bottom(s, "DATA.STREAM: VLM-ASSET-REFINERY // SYS.STATUS: ACTIVE // LATENCY: <8MS",
               "SYSTEM_MODE = AUTOMATED_LOG_PROCESSING // DATE: 2024-05-17")


def s13_agents(s):
    hud_top(s, "DATA.STREAM: AI-AGENT-ASSISTANTS // SYS.STATUS: ACTIVE // LATENCY: <5MS",
            "SYSTEM_MODE = BUILD-INTEGRATION")
    title(s, [("AI 重塑管线 02：构建智能工程超级助理", {})])
    lead(s, [("不仅是写代码，更是重构整个研发工作流。", {})], w=12.0)
    cards = [("Data Agent (Text2SQL)",
              "支持自然语言极速查询仿真大盘看板。AI 自动追踪并归因“某次版本迭代导致接管率飙升”的底层逻辑。"),
             ("Code Copilot",
              "高级辅助工具，帮助仿真工程师快速编写复杂的 C++/Python 车辆动力学插件与海量系统回归测试用例。"),
             ("Case Reviewer",
              "每日千万级里程仿真报告，由 Agent 自动阅读、汇总异常、提炼 Insight，并直接推送到企业微信。")]
    for (cl, cw), (h, b) in zip(cols(3), cards):
        panel(s, cl, 2.2, cw, 4.0, accent=CYAN)
        _, tag = add_text(s, cl + 0.25, 2.42, cw - 0.5, 0.3)
        para(tag, [("Agent Node", {})], size=10, color=CYAN, font=MONO, first=True)
        _, tf = add_text(s, cl + 0.25, 2.95, cw - 0.5, 3.1)
        para(tf, [(h, {})], size=16, color=TEXT, bold=True, first=True, space_after=12,
             line=1.15)
        para(tf, [(b, {})], size=12.5, color=MUTED, line=1.4)
    hud_bottom(s, "DATA.STREAM: AI-AGENT-ASSISTANTS // SYS.STATUS: ACTIVE // LATENCY: <5MS",
               "SYSTEM_MODE = BUILD-INTEGRATION // DATE: 2024-05-17")


def s14_philosophy(s):
    hud_top(s, "致未来的战友：拥抱“不确定”的工程哲学", "")
    _, tf = add_text(s, USABLE_L, 1.4, 12.0, 1.0)
    para(tf, [("算法团队负责跑分，我们负责修考场。", {})],
         size=34, color=TEXT, bold=True, first=True)
    _, tf2 = add_text(s, USABLE_L, 2.65, 12.0, 1.1)
    para(tf2, [("Our Mission:  ", {"color": CYAN, "bold": True, "font": MONO}),
               ("用 3DGS 和世界模型，在这个充满不确定的物理世界里，"
                "为自动驾驶打磨出一套足够“确定”的准出标准。", {})],
         size=16, color=MUTED, first=True, line=1.4)
    panel(s, USABLE_L, 3.95, USABLE_W, 2.5, fill=PANEL2, border=ORANGE)
    _, tf3 = add_text(s, USABLE_L + 0.4, 4.2, USABLE_W - 0.8, 2.0)
    para(tf3, [("Geek Principles", {})], size=13, color=ORANGE, bold=True,
         font=MONO, first=True, space_after=14)
    para(tf3, [("接口第一（Interface First）：", {"color": TEXT, "bold": True}),
               ("架构的灵活性，来自简洁、清晰的代码抽象。", {})],
         size=14.5, color=MUTED, space_after=12, line=1.35)
    para(tf3, [("不迷信日志（Data over Dogma）：", {"color": TEXT, "bold": True}),
               ("在闭环数据和真实物理面前，主观的借口往往站不住脚。", {})],
         size=14.5, color=MUTED, line=1.35)
    hud_bottom(s, "DATA.STREAM: ENGINEER-MANIFESTO // SYS.STATUS: ACTIVE // LATENCY: <1MS",
               "SYSTEM_MODE = PHILOSOPHY-DEPLOYMENT // DATE: 2024-05-20")


def s15_join(s):
    hud_top(s, "DATA.STREAM: RECRUITMENT-CALL // SYS.STATUS: ACTIVE // LATENCY: <1MS", "")
    _, tf = add_text(s, 0.8, 1.7, 11.73, 1.4, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("JOIN US.", {})], size=72, color=ORANGE, bold=True,
         align=PP_ALIGN.CENTER, first=True)
    _, tf2 = add_text(s, 1.0, 3.25, 11.33, 0.6)
    para(tf2, [("共同定义", {}), ("“虚拟未来”", {"color": ORANGE, "bold": True}),
               ("的物理法则。", {})], size=22, color=TEXT, bold=True,
         align=PP_ALIGN.CENTER, first=True)
    _, tf3 = add_text(s, 1.0, 4.0, 11.33, 0.5)
    para(tf3, [("在自动驾驶迈入 2.0 时代的浩荡浪潮中，我们缺一个你。", {})],
         size=15, color=MUTED, align=PP_ALIGN.CENTER, first=True)
    panel(s, 3.67, 4.85, 6.0, 1.2, fill=PANEL2, border=ORANGE)
    _, tf4 = add_text(s, 3.67, 4.95, 6.0, 1.0, anchor=MSO_ANCHOR.MIDDLE)
    para(tf4, [("[ 校招 / 实习转正窗口：OPEN ]", {})], size=14, color=GREEN,
         bold=True, font=MONO, align=PP_ALIGN.CENTER, first=True, space_after=8)
    para(tf4, [("Email: ", {"color": MUTED}),
               ("sim-algo-lead@xpeng.com", {"color": CYAN})],
         size=14, align=PP_ALIGN.CENTER, font=MONO)
    hud_bottom(s, "END OF DOCUMENT // SYSTEM.HALT(SUCCESS)",
               "DATA.STREAM: RECRUITMENT-CALL // SYS.STATUS: ACTIVE")


# ====================================================================
#  主驱动：PPTX 页号(1-based) -> builder
# ====================================================================
BUILDERS = {
    1:  s01_title,
    3:  s02_paradigm,
    4:  s03_dataloop,
    5:  s04_deadlock,
    6:  s05_spatiotemporal,
    7:  s06_covariate,
    8:  s07_nerf,
    9:  s08_pipeline,
    10: s09_gensim,
    11: s10_dimension,
    12: s11_topology,
    13: s12_vlm,
    14: s13_agents,
    19: s14_philosophy,
    21: s15_join,
}


def main():
    prs = Presentation(SRC)
    slides = list(prs.slides)
    for pno, builder in BUILDERS.items():
        s = slides[pno - 1]
        removed = clear_pictures(s)
        set_bg(s)
        builder(s)
        print(f"  page {pno:>2}: -{removed} pic, builder={builder.__name__}")
    prs.save(OUT)
    print("saved ->", OUT)


if __name__ == "__main__":
    main()
