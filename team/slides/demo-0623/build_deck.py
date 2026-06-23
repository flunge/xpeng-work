# -*- coding: utf-8 -*-
"""
重构【打破工程困局】整套 PPT：
 主线：仿真算法团队为 VLA 2.0 大模型“修一座追得上其进化速度的考场”。
 结构：封面 → 路线图 → [01 背景] → [02 技术演进] → [03 现状与落地] → [04 未来展望] → JOIN US
 每章前有分隔页，路线图与章节严格对应。
 复用 rebuild_editable 的底层组件 helper；屏蔽其噪声 HUD 文案（DATA.STREAM / SYSTEM_MODE / DIGITAL TWIN ...）。
 每页右上角统一 XPENG logo；封面 / 4 章节分隔 / 重建管线 / 世界模型 使用 PIL 程序化配图。
"""
import os
import rebuild_editable as R
from rebuild_editable import (
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE,
    add_box, add_text, para, panel, cols, terminal, card_block, mini_card,
    cmp_panel, simple_table, set_bg,
    BG, PANEL, PANEL2, BORDER, TEXT, MUTED, DIM, ORANGE, CYAN, GREEN, RED, LEGACY,
    MONO, CN, USABLE_L, USABLE_R, USABLE_W, L_MARGIN,
)
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
ASSET = os.path.join(HERE, "assets_gen")
AID = os.path.join(ASSET, "ai")   # gpt-image-2 生成的实景配图
OUT = os.path.join(HERE, "打破工程困局_仿真分享_30min.pptx")
WARM = RGBColor(0x1C, 0x15, 0x0E)   # 高亮卡片暖底
TOTAL = 29
DECK_FOOT = "打破工程困局 · 自动驾驶仿真算法"

# 屏蔽 rebuild_editable 内置 HUD（去除无关噪声文案）
R.hud_top = lambda *a, **k: None
R.hud_bottom = lambda *a, **k: None


# ---------------- 统一页面 chrome ----------------
def set_alpha(shape, pct):
    """给 solid 填充设置不透明度 pct(0-100)。"""
    srgb = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
    a = srgb.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))})
    srgb.append(a)


def overlay(slide, l, t, w, h, rgb=BG, opacity=45):
    sp = add_box(slide, l, t, w, h, fill=rgb)
    set_alpha(sp, opacity)
    return sp


def add_logo(slide, w=1.55):
    # 官方 XPENG 字标（由 svg_logo.py 离线栅格化）
    h = w * 150.0 / 1251.0
    slide.shapes.add_picture(os.path.join(ASSET, "logo_white.png"),
                             Inches(12.83 - w), Inches(0.34), Inches(w), Inches(h))


def ref(slide, text, t=6.72):
    _, tf = add_text(slide, USABLE_L, t, 11.0, 0.28)
    para(tf, [("◇ 参考研究  ", {"color": CYAN}), (text, {})],
         size=9, color=DIM, font=MONO, first=True)


def chrome(slide, page, section=None, dark=False):
    """通用页眉页脚：左上章节标签 + 底部页脚 + 右上 logo。"""
    if section:
        _, tf = add_text(slide, L_MARGIN, 0.32, 7.0, 0.3)
        para(tf, [(section, {})], size=10.5, color=CYAN, font=MONO, first=True)
    if not dark:
        add_box(slide, L_MARGIN, 7.05, 12.33, 0.012, fill=BORDER)
    _, lf = add_text(slide, L_MARGIN, 7.12, 7.5, 0.3)
    para(lf, [(DECK_FOOT, {})], size=9, color=DIM, font=MONO, first=True)
    _, rf = add_text(slide, 8.0, 7.12, 4.83, 0.3)
    para(rf, [(f"{page:02d}", {"color": ORANGE if not dark else CYAN}),
              (f" / {TOTAL}", {"color": DIM})],
         size=9, font=MONO, align=PP_ALIGN.RIGHT, first=True)
    add_logo(slide)


def title(slide, runs, t=0.62, size=27):
    add_box(slide, L_MARGIN, t + 0.05, 0.09, 0.5, fill=ORANGE)
    _, tf = add_text(slide, USABLE_L, t, 10.3, 0.7)
    para(tf, runs, size=size, color=TEXT, bold=True, first=True, line=1.05)


def lead(slide, runs, t=1.38, w=11.2, color=MUTED, size=14):
    _, tf = add_text(slide, USABLE_L, t, w, 0.7)
    para(tf, runs, size=size, color=color, first=True, line=1.25)


# 让被复用的 R 内部 title/lead 也指向新的（去 HUD 版式更紧凑）
R.title = title
R.lead = lead


# ---------------- 封面 ----------------
def build_cover(slide, page):
    slide.shapes.add_picture(os.path.join(AID, "cover_ai.png"),
                             Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    overlay(slide, 0, 0, 13.333, 7.5, rgb=BG, opacity=55)
    _, tf = add_text(slide, 0.85, 1.35, 11.6, 0.4)
    para(tf, [("XPENG · 自动驾驶仿真算法 / 2026 实习生技术分享", {})],
         size=12, color=CYAN, font=MONO, align=PP_ALIGN.CENTER, first=True)
    _, tt = add_text(slide, 0.85, 2.05, 11.6, 1.3, anchor=MSO_ANCHOR.MIDDLE)
    para(tt, [("打破工程困局", {})], size=60, color=TEXT, bold=True,
         align=PP_ALIGN.CENTER, first=True)
    _, ts = add_text(slide, 1.6, 3.5, 10.1, 1.0)
    para(ts, [("如何在", {}), ("“不确定”", {"color": ORANGE, "bold": True}),
              ("的物理世界里，为自动驾驶大模型构建", {}),
              ("“确定”", {"color": ORANGE, "bold": True}), ("的运行闭环", {})],
         size=17, color=TEXT, align=PP_ALIGN.CENTER, first=True, line=1.35)
    # 主线 tagline
    tl = add_box(slide, 3.43, 4.7, 6.47, 0.62, fill=PANEL2, line=ORANGE, rounded=True)
    set_alpha(tl, 70)
    _, tb = add_text(slide, 3.43, 4.72, 6.47, 0.58, anchor=MSO_ANCHOR.MIDDLE)
    para(tb, [("主线 · 修一座追得上 VLA 进化速度的考场", {})],
         size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER, first=True)
    _, sp = add_text(slide, 0.85, 5.7, 11.6, 0.4)
    para(sp, [("仿真算法负责人：", {"color": MUTED}), ("Li Kun", {"color": TEXT, "bold": True})],
         size=14, align=PP_ALIGN.CENTER, first=True)
    chrome(slide, page, dark=True)


# ---------------- 路线图（与 4 章节严格对应）----------------
def build_agenda(slide, page):
    title(slide, [("今天的路线图", {})])
    lead(slide, [("沿一条主线讲清楚：旧考场为何失灵 → 我们如何重建 → 现在做到哪、怎么服务业务 → 未来怎么自进化。", {})])
    items = [
        ("01", "背景 · 旧考场为何失灵", "先讲清自动驾驶仿真在做什么；VLA 2.0 端到端为什么让传统 LogSim / WorldSim 不再够用。", ORANGE),
        ("02", "技术演进 · 五代考场进化", "从规则 → 游戏引擎 → Gauss+Diffusion → 世界模型；以及我们的 SimWorld 重建框架。", CYAN),
        ("03", "现状 · 考场如何服务业务", "SimWorld 四条产品线如何服务量产，业务为王，Q2 可上线的硬指标。", GREEN),
        ("04", "未来 · 让考场自己进化", "VLM + Agent 重塑研发流，仿真迈向“数据工厂”，模型闭环自进化。", ORANGE),
    ]
    cw = (USABLE_W - 3 * 0.3) / 4
    for i, (num, head, desc, c) in enumerate(items):
        cl = USABLE_L + i * (cw + 0.3)
        panel(slide, cl, 2.2, cw, 4.0, accent=c)
        _, tf = add_text(slide, cl + 0.24, 2.45, cw - 0.48, 3.5)
        para(tf, [(num, {})], size=44, color=c, bold=True, font=MONO, first=True,
             space_after=6)
        para(tf, [(head, {})], size=15, color=TEXT, bold=True, space_after=12, line=1.2)
        para(tf, [(desc, {})], size=12, color=MUTED, line=1.4)
    chrome(slide, page, section="ROADMAP · 路线图")


# ---------------- 章节分隔页 ----------------
DIV_META = {
    1: ("01", "背景", "旧考场为何失灵", "BACKGROUND · WHY", ORANGE),
    2: ("02", "技术演进", "五代考场的进化之路", "EVOLUTION · HOW", CYAN),
    3: ("03", "现状与落地", "今天的考场与战果", "STATUS · WHAT", GREEN),
    4: ("04", "未来展望", "让考场自己进化", "FUTURE · NEXT", ORANGE),
}


def build_divider(slide, idx, page):
    num, zh, sub, en, c = DIV_META[idx]
    aimg = {1: "div1_crisis", 2: "div2_evolution", 3: "div3_framework", 4: "div4_future"}[idx]
    slide.shapes.add_picture(os.path.join(AID, f"{aimg}.png"),
                             Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    overlay(slide, 0, 0, 13.333, 7.5, rgb=BG, opacity=48)
    overlay(slide, 0, 0, 7.6, 7.5, rgb=BG, opacity=42)   # 左侧再压暗，保证标题可读
    add_box(slide, 0, 0, 0.16, 7.5, fill=c)
    _, ef = add_text(slide, 1.1, 2.05, 10.0, 0.4)
    para(ef, [(en, {})], size=14, color=c, font=MONO, first=True)
    _, nf = add_text(slide, 1.05, 2.45, 4.0, 1.6)
    para(nf, [(num, {})], size=110, color=c, bold=True, font=MONO, first=True)
    _, tf = add_text(slide, 3.7, 2.7, 8.5, 1.1)
    para(tf, [(zh, {})], size=46, color=TEXT, bold=True, first=True)
    _, sf = add_text(slide, 3.75, 3.95, 8.5, 0.6)
    para(sf, [(sub, {})], size=20, color=MUTED, first=True)
    chrome(slide, page, dark=True)


# ---------------- 技术演进里程碑 G1→G5 ----------------
def build_evolution(slide, page):
    title(slide, [("五代仿真进化：从“纯规则”到“AI 造世界”", {})])
    lead(slide, [("仿真大致走过五个世代，每一代都在往两个方向使劲：真实度更高，交互性更强。", {})])
    gens = [("G1", "规则仿真", ["纯规则 / 脚本驱动", "NPC 走 If-Else", "真实感差·长尾靠手写"], LEGACY, False),
            ("G2", "游戏引擎", ["UE / CARLA 搭场景", "可交互·可闭环", "Sim2Real gap 大·画面假"], LEGACY, False),
            ("G3", "Gauss + Diffusion", ["3DGS 重建 + 扩散修复", "照片级·可编辑·实时", "真实与交互兼得"], CYAN, True),
            ("G4", "世界模型", ["生成式·学习物理常识", "想象未知长尾场景", "视频 / DiT 生成"], ORANGE, False),
            ("G5", "VLA + RL + Agent", ["闭环自进化", "AI 自动造考场", "重构研发流"], ORANGE, False)]
    cs = cols(5, gap=0.22)
    add_box(slide, USABLE_L, 2.62, USABLE_W, 0.02, fill=BORDER)
    for (cl, cw), (g, name, bullets, c, cur) in zip(cs, gens):
        add_box(slide, cl + cw / 2 - 0.06, 2.55, 0.12, 0.12, fill=c, rounded=True)
        panel(slide, cl, 2.85, cw, 3.2, accent=c, fill=(WARM if cur else PANEL))
        _, tf = add_text(slide, cl + 0.16, 3.04, cw - 0.3, 2.9)
        para(tf, [(g, {})], size=24, color=c, bold=True, font=MONO, first=True)
        para(tf, [(name, {})], size=13, color=TEXT, bold=True, space_after=8, line=1.1)
        for b in bullets:
            para(tf, [("· " + b, {})], size=10.3, color=MUTED, space_after=4, line=1.12)
    # 当前跃迁标记（G3 当前，探索 G4）
    bl = cs[2][0]
    bw = cs[3][0] + cs[3][1] - cs[2][0]
    add_box(slide, bl, 6.2, bw, 0.5, fill=WARM, line=ORANGE, rounded=True)
    _, bf = add_text(slide, bl, 6.22, bw, 0.46, anchor=MSO_ANCHOR.MIDDLE)
    para(bf, [("我们目前处在 G3（Gauss + Diffusion），同时在往 G4 世界模型探索", {})],
         size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER, first=True)
    chrome(slide, page, section="02 · 技术演进")


# ---------------- 现状：团队拓扑 ----------------
def build_team(slide, page):
    title(slide, [("我们在做什么：四条产品线", {})])
    lead(slide, [("团队在上海，分四条产品线，核心是把前沿算法真正落到量产链路上——业务定义“考什么”，技术定义“怎么考”。", {})])
    lines = [("场景 & 生产", "把真实路采重建成可直接用的仿真场景。", "极速重建 · 场景编辑 · 闭环场景集", CYAN),
             ("SIL 软件在环", "纯软件环境批量跑模型回归测试。", "车型泛化 · 画质质检 · 长里程验证", CYAN),
             ("HIL 硬件在环", "接入真实域控/台架做硬件在环仿真。", "链路部署 · 高保真慢速回放", ORANGE),
             ("Agents 智能体", "用 Agent 自动化研发里的重复劳动。", "复现率分析 · 自动出图 · Prompt 工程", ORANGE)]
    cw = (USABLE_W - 3 * 0.3) / 4
    for i, (head, plain, terms, c) in enumerate(lines):
        cl = USABLE_L + i * (cw + 0.3)
        panel(slide, cl, 2.15, cw, 2.55, accent=c)
        _, tf = add_text(slide, cl + 0.2, 2.4, cw - 0.4, 2.1)
        para(tf, [(f"产品线 0{i+1}", {})], size=9.5, color=DIM, font=MONO, first=True,
             space_after=4)
        para(tf, [(head, {})], size=15, color=TEXT, bold=True, space_after=8, line=1.15)
        para(tf, [(plain, {})], size=12, color=MUTED, space_after=8, line=1.25)
        para(tf, [(terms, {})], size=9.5, color=DIM, line=1.2)
    panel(slide, USABLE_L, 5.0, USABLE_W, 1.45, fill=PANEL2, border=ORANGE)
    _, tf = add_text(slide, USABLE_L + 0.3, 5.2, USABLE_W - 0.6, 1.05, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("分工逻辑：", {"color": CYAN, "bold": True}),
              ("业务 Owner 定义“考什么”，技术 Owner 定义“怎么考”。", {})],
         size=14, color=TEXT, first=True, space_after=8, line=1.3)
    para(tf, [("仿真不只是辅助工具，在端到端时代它更像上车前的", {}),
              ("准入门槛（Gating）", {"color": ORANGE, "bold": True}),
              ("——模型能不能上车，得先过我们这一关。", {})],
         size=14, color=TEXT, line=1.3)
    chrome(slide, page, section="03 · 现状与落地")


# ---------------- 现状：Q2 战报 ----------------
def build_battle(slide, page):
    title(slide, [("业务落地数据：前沿不是 PPT，是能上线的数字", {})])
    lead(slide, [("下面是 Q2 的实测数据，每个数字背后，都是给量产省下的成本、追回的时间、兜住的风险。", {})])
    metrics = [("91 min", "单场景重建时间", "提速 2.6×，考场当天建、当天用", CYAN),
               ("1000 km/天", "HIL 台架仿真里程", "效率 1:2.5，用仿真替真车路测", CYAN),
               ("32 dB+", "场景重建保真度 PSNR", "接近实拍，模型可直接训练", GREEN),
               ("90%", "坏场景质检召回率", "自动拦截劣质数据，省人工筛查", GREEN),
               ("74%", "线上问题闭环复现率", "7 成问题仿真可复现定位，冲 80%", ORANGE),
               ("4 / 12", "Agent 专项已上线", "变道/限速/红灯/居中，自动验收", ORANGE)]
    cw = (USABLE_W - 2 * 0.3) / 3
    ch = 1.95
    for i, (big, label, sub, c) in enumerate(metrics):
        r, col = divmod(i, 3)
        cl = USABLE_L + col * (cw + 0.3)
        ct = 2.2 + r * (ch + 0.3)
        panel(slide, cl, ct, cw, ch, accent=c)
        _, tf = add_text(slide, cl + 0.25, ct + 0.22, cw - 0.5, ch - 0.4)
        para(tf, [(big, {})], size=34, color=c, bold=True, font=MONO, first=True,
             space_after=4)
        para(tf, [(label, {})], size=13, color=TEXT, bold=True, space_after=4, line=1.15)
        para(tf, [(sub, {})], size=11, color=MUTED, line=1.2)
    chrome(slide, page, section="03 · 现状与落地")


# ---------------- 未来：路线规划 NOW/NEXT/FUTURE ----------------
def build_roadmap(slide, page):
    title(slide, [("未来规划：从“修考场”到“造世界”", {})])
    lead(slide, [("分三步走，让仿真慢慢变成模型自我迭代的引擎。", {})])
    phases = [("NOW", "夯实考场", ["闭环复现率 → 80%+", "SIL 效率 → 1:25",
                                "HIL 1000km/天 稳定 95%+", "车型泛化覆盖全系"], CYAN),
              ("NEXT", "智能化管线", ["Agent 自动复现/归因/修复", "VLM 资产沉淀闭环",
                                  "12 专项 → 全量上线", "建立 Agent 自迭代机制"], GREEN),
              ("FUTURE", "生成式世界", ["世界模型生成长尾场景", "DiT × 3DGS 混合落地",
                                    "VLA + RL 闭环自进化", "仿真即“数据工厂”"], ORANGE)]
    cw = (USABLE_W - 2 * 0.55) / 3
    for i, (ph, head, bullets, c) in enumerate(phases):
        cl = USABLE_L + i * (cw + 0.55)
        panel(slide, cl, 2.2, cw, 4.1, accent=c)
        _, tf = add_text(slide, cl + 0.28, 2.45, cw - 0.56, 3.6)
        para(tf, [(ph, {})], size=22, color=c, bold=True, font=MONO, first=True)
        para(tf, [(head, {})], size=16, color=TEXT, bold=True, space_after=12, line=1.15)
        for b in bullets:
            para(tf, [("• ", {"color": c}), (b, {})], size=12.5, color=MUTED,
                 space_after=9, line=1.25)
        if i < 2:
            _, ar = add_text(slide, cl + cw + 0.08, 3.9, 0.42, 0.6)
            para(ar, [("→", {})], size=22, color=ORANGE, align=PP_ALIGN.CENTER, first=True)
    chrome(slide, page, section="04 · 未来展望")


# ---------------- 未来：三条箴言 ----------------
def build_mindset(slide, page):
    title(slide, [("给在座同学的三点建议", {})])
    lead(slide, [("和在座的硕博同学聊聊：怎么在“不确定”里做出“确定”。", {})])
    cards = [("不确定 → 确定",
              "大模型的本质是概率，工程的本质是可复现。用闭环、用数据、用裁判机制，把“大概率对”变成“可验证的确定”。", CYAN),
             ("论文 → 产线",
              "能跑出 demo 不算赢，能在量产线稳定出数字才算。再前沿的算法，落不到 PSNR、复现率、效率上，就只是漂亮的 PPT。", GREEN),
             ("单点 → 闭环",
              "不要迷恋单个算法的 SOTA，要建会自己进化的系统。数据飞轮 + Agent 自动化，才是 Scaling Law 真正的工程含义。", ORANGE)]
    for (cl, cw), (head, body, c) in zip(cols(3), cards):
        panel(slide, cl, 2.25, cw, 3.9, accent=c)
        _, tf = add_text(slide, cl + 0.26, 2.5, cw - 0.52, 3.4)
        para(tf, [(head, {})], size=18, color=c, bold=True, first=True, space_after=14,
             line=1.15)
        para(tf, [(body, {})], size=13, color=MUTED, line=1.45)
    chrome(slide, page, section="04 · 未来展望")


# ---------------- 技术演进：DrivingGaussian 管线（含点云配图）----------------
def build_pipeline_img(slide, page):
    title(slide, [("SimWorld 重建管线：从真实路采到可编辑 3DGS", {})])
    # 右侧点云配图
    fx, fw = 8.05, 4.78
    fh = fw * 9 / 16
    slide.shapes.add_picture(os.path.join(AID, "street_splat.png"),
                             Inches(fx), Inches(1.7), Inches(fw), Inches(fh))
    add_box(slide, fx, 1.7, fw, fh, line=BORDER, line_w=1)
    _, cap = add_text(slide, fx + 0.1, 1.7 + fh + 0.06, fw - 0.2, 0.3)
    para(cap, [("场景图 3DGS：把真实街景重建为可编辑点云世界", {})],
         size=10, color=CYAN, font=MONO, first=True)
    # 左侧 4 步纵向管线（SimWorld 真实流程）
    steps = [("① 数据预处理", "去畸变 + 语义分割 + 位姿优化（DPVO），相机/LiDAR 时空对齐。", CYAN, False),
             ("② 几何构建", "COLMAP SfM / LiDAR 点云 / MVSNet 深度 → 稠密点云 + 深度图。", CYAN, False),
             ("③ 场景图 3DGS 训练", "背景/地面/车辆/行人/天空分离建模（Street Gaussians · Reconic）。", ORANGE, True),
             ("④ 修复 + 仿真部署", "Difix/NVFixer 扩散修复伪影 → 仿真接口新视角渲染 + 重畸变。", CYAN, False)]
    for i, (head, body, c, key) in enumerate(steps):
        ct = 1.6 + i * 0.92
        card_block(slide, USABLE_L, ct, 7.0, 0.8, head, body, accent=c,
                   head_size=12.5, body_size=10.5, head_color=(ORANGE if key else TEXT))
        if i < 3:
            _, ar = add_text(slide, USABLE_L + 0.2, ct + 0.78, 0.4, 0.2)
            para(ar, [("▼", {})], size=9, color=DIM, first=True)
    panel(slide, USABLE_L, 5.35, USABLE_W, 1.1, fill=PANEL2, border=ORANGE)
    _, tf = add_text(slide, USABLE_L + 0.3, 5.52, USABLE_W - 0.6, 0.8, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("双 Pipeline：", {"color": ORANGE, "bold": True}),
              ("LiDAR 主力（点云高精）+ Vision 纯视觉（MVSNet 深度），覆盖有/无激光雷达车型；前馈式 EvoSplat 进一步免去逐场景优化。", {})],
         size=13, color=TEXT, first=True, line=1.3)
    chrome(slide, page, section="02 · 技术演进")


# ---------------- 技术演进：Gen-Sim 世界模型（含世界模型配图）----------------
def build_worldmodel_img(slide, page):
    title(slide, [("G4 · 世界模型：让环境会“预测未来”", {})])
    lead(slide, [("环境不再是被动的贴图，而是有物理常识、会自己反应的智能体——主要用来解决“仿真太假、NPC 太呆”这两个老问题。", {})])
    fx, fw = 0.72, 4.6
    slide.shapes.add_picture(os.path.join(AID, "world_model.png"),
                             Inches(fx), Inches(2.05), Inches(fw), Inches(3.5))
    add_box(slide, fx, 2.05, fw, 3.5, line=BORDER, line_w=1)
    _, cap = add_text(slide, fx + 0.1, 5.18, fw - 0.2, 0.3)
    para(cap, [("Digital World Model · 生成式时空推演", {})],
         size=10, color=CYAN, font=MONO, first=True)
    cards = [("会“脑补”未来", "给它当前画面 + 本车要做的动作，它能推演出接下来几秒整段场景怎么变。", CYAN),
             ("会“讲道理”", "不再是死板 If-Else：“我一急刹，后车按常理减速避让”——像真人一样博弈。", ORANGE),
             ("会“变天”", "同一个路口，一键生成雨天、雾天、黄昏等各种天气光照，长尾场景随取随用。", GREEN)]
    bx = 5.6
    for i, (head, body, c) in enumerate(cards):
        card_block(slide, bx, 2.05 + i * 1.2, USABLE_R - bx, 1.05, head, body,
                   accent=c, head_size=14, body_size=12)
    chrome(slide, page, section="02 · 技术演进")


# ---------------- 新增：含配图的内容页（替换枯燥文本页）----------------
def build_dataloop_img(slide, page):
    title(slide, [("数据闭环：VLA 模型进化的", {}), ("核心动力", {"color": ORANGE})])
    fx, fw = 0.72, 3.85
    slide.shapes.add_picture(os.path.join(AID, "flywheel.png"),
                             Inches(fx), Inches(1.7), Inches(fw), Inches(fw * 820 / 860))
    add_box(slide, fx, 1.7, fw, fw * 820 / 860, line=BORDER, line_w=1)
    _, cap = add_text(slide, fx, 5.55, fw, 0.3)
    para(cap, [("Data Flywheel · 仿真驱动的强化学习闭环", {})],
         size=10, color=CYAN, font=MONO, align=PP_ALIGN.CENTER, first=True)
    bx = 4.95
    _, tf = add_text(slide, bx, 1.7, USABLE_R - bx, 1.0)
    para(tf, [("仿真不只是测试工具，对车端 VLA 2.0 来说，它还是做", {}),
              ("强化学习（RL）", {"color": CYAN, "bold": True}),
              ("时的虚拟教练。", {})], size=15, color=TEXT, first=True, line=1.3)
    fw3 = [("① 海量“模拟考”", "每天在仿真里跑千万级里程，等于给新版本做海量模拟考——快速判断它能不能上车。", CYAN),
           ("② 专挑软肋反复练", "把路测捞到的疑难场景自动泛化成上百个变体，专打模型最容易出错的地方。", ORANGE),
           ("③ 择优留强汰弱", "对比各版本仿真表现，留下真正变强的、淘汰“刷分却变笨”的，保证迭代不跑偏。", GREEN)]
    for i, (h, b, c) in enumerate(fw3):
        card_block(slide, bx, 2.75 + i * 1.12, USABLE_R - bx, 1.0, h, b,
                   accent=c, head_size=13, body_size=11)
    chrome(slide, page, section="01 · 背景")


def build_nerf_img(slide, page):
    title(slide, [("神经场景基石：从隐式 NeRF 到显式 3DGS 的质变", {})])
    slide.shapes.add_picture(os.path.join(AID, "nerf_3dgs.png"),
                             Inches(0.72), Inches(1.55), Inches(6.2), Inches(3.43))
    add_box(slide, 0.72, 1.55, 6.2, 3.43, line=BORDER, line_w=1)
    add_box(slide, 0.72, 1.55, 3.1, 0.34, fill=PANEL2)
    add_box(slide, 3.82, 1.55, 3.1, 0.34, fill=PANEL2)
    _, l1 = add_text(slide, 0.82, 1.57, 2.9, 0.3, anchor=MSO_ANCHOR.MIDDLE)
    para(l1, [("NeRF · 隐式 MLP（模糊/慢）", {})], size=9.5, color=MUTED, font=MONO, first=True)
    _, l2 = add_text(slide, 3.92, 1.57, 2.9, 0.3, anchor=MSO_ANCHOR.MIDDLE)
    para(l2, [("3DGS · 显式点云（清晰/极速）", {})], size=9.5, color=CYAN, font=MONO, first=True)
    rows = [("存储成本", "MLP ~50MB（极小）", "Splats ~300MB（较大）"),
            ("渲染延迟", "1–2s / 帧（极慢）", "100+ FPS（极速）"),
            ("编辑性", "极难（Implicit Field）", "易编辑（Point-based）")]
    bx = 7.25
    _, tf = add_text(slide, bx, 1.55, USABLE_R - bx, 0.4)
    para(tf, [("关键对比", {})], size=12, color=ORANGE, bold=True, font=MONO, first=True)
    for i, (k, a, b) in enumerate(rows):
        ct = 2.0 + i * 1.0
        panel(slide, bx, ct, USABLE_R - bx, 0.88, fill=PANEL)
        _, c = add_text(slide, bx + 0.2, ct + 0.1, USABLE_R - bx - 0.4, 0.7)
        para(c, [(k, {})], size=11, color=TEXT, bold=True, first=True, space_after=3)
        para(c, [("NeRF ", {"color": DIM}), (a, {"color": LEGACY}),
                 ("   →   ", {"color": DIM}), (b, {"color": CYAN})],
             size=10.5, line=1.1)
    _, tk = add_text(slide, 0.72, 5.15, 6.2, 0.8)
    para(tk, [("3DGS 以可接受的存储代价，换来两个数量级的渲染提速与可编辑性——支撑实时闭环仿真的关键基石。", {})],
         size=12, color=MUTED, first=True, line=1.3)
    ref(slide, "NeRF (Mildenhall, ECCV’20) · 3D Gaussian Splatting (Kerbl, SIGGRAPH’23)")
    chrome(slide, page, section="02 · 技术演进")


def build_drift_img(slide, page):
    title(slide, [("深度推演：协变量偏移（Covariate Shift）灾难链条", {})])
    lead(slide, [("本车 Action 与录像产生 0.1° 偏差 → 5s 后偏离 1.5m → 空间错乱 → 模型崩溃。", {})])
    fx, fw = 0.72, 6.2
    fh = fw * 9 / 16
    slide.shapes.add_picture(os.path.join(AID, "drift.png"),
                             Inches(fx), Inches(2.0), Inches(fw), Inches(fh))
    add_box(slide, fx, 2.0, fw, fh, line=BORDER, line_w=1)
    _, cap = add_text(slide, fx + 0.1, 2.0 + fh + 0.06, fw - 0.2, 0.3)
    para(cap, [("青线=录像参考轨迹  橙线=本车真实轨迹（发散→崩溃）", {})],
         size=9.5, color=MUTED, font=MONO, first=True)
    bx = 7.25
    cards = [("黑盒测试的唯一出口", "VLA 内部无 BBox 接口，只能通过闭环（Closed-loop）观察——不接管，车会不会撞？", ORANGE),
             ("仿真底线", "必须提供实时渲染，以毫秒级响应 VLA 的每一个转角输出，否则无法测试偏差后的纠偏能力。", CYAN)]
    for i, (h, b, c) in enumerate(cards):
        card_block(slide, bx, 2.0 + i * 1.75, USABLE_R - bx, 1.55, h, b,
                   accent=c, head_size=15, body_size=12.5)
    chrome(slide, page, section="02 · 技术演进")


def build_vlm_img(slide, page):
    title(slide, [("AI 重塑管线 01：VLM 自动化资产沉淀", {})])
    lead(slide, [("现实痛点：", {"color": CYAN, "bold": True}),
                 ("每天产生 TB 级的路测 Log 录像，靠人力根本看不过来。", {})])
    slide.shapes.add_picture(os.path.join(AID, "vlm_scan.png"),
                             Inches(0.72), Inches(2.0), Inches(6.2), Inches(3.43))
    add_box(slide, 0.72, 2.0, 6.2, 3.43, line=BORDER, line_w=1)
    _, cap = add_text(slide, 0.82, 5.05, 6.0, 0.3)
    para(cap, [("VLM Scanner · 自动语义标注与缺陷检测", {})],
         size=10, color=CYAN, font=MONO, first=True)
    bx = 7.25
    cards = [("场景理解 (Semantic)", "自动扫描输出“无保护左转”“复杂路口博弈”等高价值语义标签。", CYAN),
             ("缺陷发现 (Defect)", "自动寻找 VLA 表现不佳（如急刹车）的薄弱片段，直接送入飞轮重训。", ORANGE),
             ("资产沉淀 (Refinery)", "TB 级 Log 自动精炼为高优先级训练集，把“人看不过来”变成“机器筛得过来”。", GREEN)]
    for i, (h, b, c) in enumerate(cards):
        card_block(slide, bx, 2.0 + i * 1.18, USABLE_R - bx, 1.02, h, b,
                   accent=c, head_size=13.5, body_size=11.5)
    chrome(slide, page, section="04 · 未来展望")


def build_agents_img(slide, page):
    title(slide, [("AI 重塑管线 02：Agent 现在已经在替我们做这些事", {})])
    lead(slide, [("不只是写代码，Agent 已经进入我们仿真研发的日常工作流。", {})])
    agents = [("Coding Agent", "代码库改造 + 工作流自动化：重构历史代码、自动写脚手架 / 测试 / CI。", CYAN),
              ("Case 打标 Agent", "自动给海量路测、仿真 case 打语义标签，替代人工标注、沉淀场景库。", ORANGE),
              ("复现率分析 Agent", "自动判断线上问题能否在仿真复现、量化复现率并定位回放差异。", GREEN),
              ("Metric 分析 Agent", "自动读仿真大盘指标，做异常归因、生成结论日报。", CYAN),
              ("Smart Planner", "自然语言 → 自动编排仿真任务与实验计划，免手工配 pipeline。", ORANGE),
              ("Data Agent (Text2SQL)", "用大白话查大盘——“为什么这版接管率飙升”，一问即答。", GREEN)]
    cw = (USABLE_W - 2 * 0.3) / 3
    for i, (h, b, c) in enumerate(agents):
        r, col = divmod(i, 3)
        cl = USABLE_L + col * (cw + 0.3)
        ct = 2.2 + r * 2.15
        card_block(slide, cl, ct, cw, 1.95, h, b, accent=c, head_size=13.5, body_size=11.5)
    chrome(slide, page, section="04 · 未来展望")


# ---------------- 未来：Agent 还能为我们做什么 ----------------
def build_agent_future(slide, page):
    title(slide, [("未来：Agent 还能为我们做什么", {})])
    lead(slide, [("顺着现在的自动化往前看，Agent 会从“帮手”慢慢变成一支能独立干活的“工程团队”。", {})])
    fx, fw = 0.72, 5.3
    fh = fw * 9 / 16
    slide.shapes.add_picture(os.path.join(AID, "agent_team.png"),
                             Inches(fx), Inches(2.05), Inches(fw), Inches(fh))
    add_box(slide, fx, 2.05, fw, fh, line=BORDER, line_w=1)
    _, cap = add_text(slide, fx + 0.1, 2.05 + fh + 0.06, fw - 0.2, 0.3)
    para(cap, [("多 Agent 协同：Planner · Coder · Reviewer · Tester", {})],
         size=10, color=CYAN, font=MONO, first=True)
    bx = 6.35
    cards = [("从“发现问题”到“修好问题”", "Agent 直接定位代码/数据根因、提交修复 MR，人只做 Review。", CYAN),
             ("多 Agent 协同流水线", "Planner→Coder→Reviewer→Tester 分工协作，端到端跑完一个需求。", ORANGE),
             ("自主出“考题课程”", "按模型当前弱点，自动生成由易到难的场景课程，驱动模型自进化。", GREEN),
             ("一句话造世界", "结合世界模型，自然语言即可生成符合物理的长尾场景。", CYAN)]
    for i, (h, b, c) in enumerate(cards):
        card_block(slide, bx, 2.05 + i * 1.08, USABLE_R - bx, 0.95, h, b,
                   accent=c, head_size=13, body_size=11)
    _, nt = add_text(slide, USABLE_L, 6.35, 11.6, 0.4)
    para(nt, [("再往后：", {"color": ORANGE, "bold": True}),
              ("仿真团队的角色，会从“修考场”变成“带一支 AI 工程团队”，让研发流程自己迭代自己。", {})],
         size=12.5, color=TEXT, first=True)
    chrome(slide, page, section="04 · 未来展望")


# ---------------- 背景：双重死局（含配图）----------------
def build_deadlock_img(slide, page):
    title(slide, [("传统仿真架构面对 VLA 的", {}), ("两难困境", {"color": ORANGE})])
    p = panel(slide, USABLE_L, 1.5, USABLE_W, 0.5, fill=PANEL2)
    _, tf = add_text(slide, USABLE_L + 0.2, 1.5, USABLE_W - 0.4, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("必须找到一种 ", {}),
              ("“既有 LogSim 真实度，又有 WorldSim 交互性”", {"color": CYAN, "bold": True}),
              (" 的新方案。", {})], size=13.5, color=TEXT, align=PP_ALIGN.CENTER, first=True)
    fx, fw = 0.72, 5.5
    fh = fw * 9 / 16
    slide.shapes.add_picture(os.path.join(AID, "deadlock.png"),
                             Inches(fx), Inches(2.3), Inches(fw), Inches(fh))
    add_box(slide, fx, 2.3, fw, fh, line=BORDER, line_w=1)
    _, cap = add_text(slide, fx + 0.1, 2.3 + fh + 0.06, fw - 0.2, 0.3)
    para(cap, [("LogSim 录像回放  vs  WorldSim 白模规则 —— 两条路各有各的死角", {})],
         size=9.5, color=MUTED, font=MONO, first=True)
    bx = 6.5
    cards = [("LogSim（数据回放）", "开环回放实车 Log，环境不随本车动作改变；VLA 无 BBox 接口，无法 Assert。", ORANGE),
             ("WorldSim（白模规则）", "手工 3D 资产 + NPC 规则；Sim2Real Gap 巨大，视觉大模型不认“游戏画面”。", ORANGE)]
    for i, (h, b, c) in enumerate(cards):
        ct = 2.3 + i * 2.0
        panel(slide, bx, ct, USABLE_R - bx, 1.75, accent=RED)
        _, t = add_text(slide, bx + 0.25, ct + 0.2, USABLE_R - bx - 0.5, 1.4)
        para(t, [(h, {})], size=15, color=TEXT, bold=True, first=True, space_after=8)
        para(t, [(b, {})], size=12, color=MUTED, line=1.35, space_after=8)
        para(t, [("DEAD END · VLA INCOMPATIBLE", {})], size=10, color=RED, bold=True, font=MONO)
    chrome(slide, page, section="01 · 背景")


# ---------------- 技术演进：时空连续性编辑（含配图）----------------
def build_spatio_img(slide, page):
    title(slide, [("时空连续性编辑：从", {}), ("“录像重放”", {"color": ORANGE}),
                  ("到", {}), ("“无限变体考场”", {"color": ORANGE})])
    lead(slide, [("自研 3DGS 引擎 + 时空连续性算法，对场景中每个 Gaussian 实体实时干预，彻底解决 2D 扩散模型的视角闪烁与畸变。", {})])
    fx, fw = 0.72, 5.6
    fh = fw * 9 / 16
    slide.shapes.add_picture(os.path.join(AID, "spatiotemporal.png"),
                             Inches(fx), Inches(2.05), Inches(fw), Inches(fh))
    add_box(slide, fx, 2.05, fw, fh, line=BORDER, line_w=1)
    _, cap = add_text(slide, fx + 0.1, 2.05 + fh + 0.06, fw - 0.2, 0.3)
    para(cap, [("从 Baseline Log 派生无限高保真 Corner Case", {})],
         size=9.5, color=CYAN, font=MONO, first=True)
    bx = 7.25
    cards = [("NPC 变轨 (Trajectory Hack)", "抽离原始 Log 车辆并改轨迹，令其对本车“强行 Cut-in”。", CYAN),
             ("资产替换 (Asset Swap)", "像素级抠除路边的树，无缝替换为横穿马路的行人。", ORANGE),
             ("无限繁衍 (Infinite Spawning)", "同一地段瞬间生成上百个差异化高保真 Corner Case。", GREEN)]
    for i, (h, b, c) in enumerate(cards):
        card_block(slide, bx, 2.05 + i * 1.18, USABLE_R - bx, 1.02, h, b,
                   accent=c, head_size=13.5, body_size=11.5)
    chrome(slide, page, section="02 · 技术演进")


# ---------------- 技术演进：世界模型·2026 前沿思辨 ----------------
def build_wm_frontier(slide, page):
    title(slide, [("G4 · 世界模型 2026 前沿思辨与我们的位置", {})])
    lead(slide, [("2026 年世界模型很热，但业界对“世界模型到底是什么”分成两派——这也正好帮我们想清楚自己做仿真的位置。", {})])
    fx, fw = 0.72, 5.6
    fh = fw * 9 / 16
    slide.shapes.add_picture(os.path.join(AID, "wm_frontier.png"),
                             Inches(fx), Inches(2.05), Inches(fw), Inches(fh))
    add_box(slide, fx, 2.05, fw, fh, line=BORDER, line_w=1)
    _, cap = add_text(slide, fx + 0.1, 2.05 + fh + 0.06, fw - 0.2, 0.3)
    para(cap, [("理解（Understanding） vs 生成（Generation）两条技术路线", {})],
         size=9.5, color=CYAN, font=MONO, first=True)
    bx = 7.25
    card_block(slide, bx, 2.05, USABLE_R - bx, 1.55,
               "理解派 · Understanding",
               "学环境内部表征与物理常识，服务规划推理。代表：LeCun 离开 Meta 创立 AMI Labs，押注 JEPA / V-JEPA（预测潜空间而非像素）。",
               accent=CYAN, head_size=14, body_size=11.5)
    card_block(slide, bx, 3.75, USABLE_R - bx, 1.55,
               "生成派 · Generation",
               "可交互、物理可信的视频 / 3D 世界生成器。代表：DeepMind Genie 3（实时交互 3D 世界）、Fei-Fei Li 的 World Labs（3D 空间智能）。",
               accent=ORANGE, head_size=14, body_size=11.5)
    panel(slide, USABLE_L, 5.4, USABLE_W, 1.05, fill=PANEL2, border=ORANGE)
    _, tf = add_text(slide, USABLE_L + 0.25, 5.5, USABLE_W - 0.5, 0.85)
    para(tf, [("争议点：", {"color": ORANGE, "bold": True}),
              ("能生成逼真画面 ≠ 真正理解物理（长时一致性、隐藏状态、几何一致，都还是难题）。", {})],
         size=11.5, color=TEXT, first=True, space_after=4, line=1.2)
    para(tf, [("我们的位置：", {"color": GREEN, "bold": True}),
              ("不追 AGI 级通用理解，做 3D-grounded 生成（3DGS × DiT）——环视几何一致 + 可编辑 + 闭环可控，把世界模型变成“测得准”的考场。", {})],
         size=11.5, color=TEXT, line=1.2)
    ref(slide, "综述《Understanding World or Predicting Future?》(ACM CSUR’25) · AMI Labs · World Labs · Genie 3", t=6.55)
    chrome(slide, page, section="02 · 技术演进")


# ---------------- 技术演进：世界模型对仿真考场的启示 ----------------
def build_wm_grounded(slide, page):
    title(slide, [("对仿真考场的启示：要“测得准”，不只是“看起来真”", {})])
    lead(slide, [("自动驾驶不需要 AGI 级通用理解，需要的是对物理足够确定、可控、几何一致的生成考场。", {})])
    fx, fw = 0.72, 5.6
    fh = fw * 9 / 16
    slide.shapes.add_picture(os.path.join(AID, "wm_grounded.png"),
                             Inches(fx), Inches(2.05), Inches(fw), Inches(fh))
    add_box(slide, fx, 2.05, fw, fh, line=BORDER, line_w=1)
    _, cap = add_text(slide, fx + 0.1, 2.05 + fh + 0.06, fw - 0.2, 0.3)
    para(cap, [("3D-grounded 生成：环视多相机几何天然一致", {})],
         size=9.5, color=CYAN, font=MONO, first=True)
    bx = 7.25
    cards = [("纯 2D 视频路线的软肋", "Sora / Genie 式 2D 生成：多摄像头几何不一致、长时漂移、隐藏状态丢失——恰是自动驾驶命门。", ORANGE),
             ("我们的解法 · 3D-grounded", "在 3D 特征空间生成（3DGS × DiT）：环视天然几何一致 + 显式可编辑 + 闭环可控。", CYAN),
             ("落点", "把“世界模型”从“逼真演示”变成“可量产的确定性考场”——对齐 PSNR / 复现率 / 闭环硬指标。", GREEN)]
    for i, (h, b, c) in enumerate(cards):
        card_block(slide, bx, 2.05 + i * 1.18, USABLE_R - bx, 1.02, h, b,
                   accent=c, head_size=13.5, body_size=11.5)
    chrome(slide, page, section="02 · 技术演进")


# ---------------- 背景：自动驾驶仿真 101 ----------------
def build_sim_intro(slide, page):
    title(slide, [("先说清楚：自动驾驶仿真在为谁、造什么“考场”", {})])
    lead(slide, [("在虚拟世界里跑里程、把各种极端 case 都撞一遍，让模型“该犯的错先在仿真里犯完”——比真车路测更安全、更便宜、也更快。", {})])
    fx, fw = 0.72, 5.6
    fh = fw * 9 / 16
    slide.shapes.add_picture(os.path.join(AID, "sim_intro.png"),
                             Inches(fx), Inches(2.05), Inches(fw), Inches(fh))
    add_box(slide, fx, 2.05, fw, fh, line=BORDER, line_w=1)
    _, cap = add_text(slide, fx + 0.1, 2.05 + fh + 0.06, fw - 0.2, 0.3)
    para(cap, [("虚拟考场：可规模化、可重复、零风险地验证自动驾驶模型", {})],
         size=9.5, color=CYAN, font=MONO, first=True)
    bx = 7.25
    cards = [("为什么要仿真？", "真车路测贵、慢、危险，长尾场景难复现；虚拟考场能低成本、规模化、可重复地把模型“考”到位。", CYAN),
             ("LogSim · 数据回放", "把真实路测录像原样重放——“真”，但“死”：开环、环境不随本车动作改变。", ORANGE),
             ("WorldSim · 虚拟搭建", "用游戏引擎 / 白模手工搭场景——“活”，但“假”：可交互，但 Sim2Real 画面失真。", GREEN)]
    for i, (h, b, c) in enumerate(cards):
        card_block(slide, bx, 2.05 + i * 1.18, USABLE_R - bx, 1.02, h, b,
                   accent=c, head_size=13.5, body_size=11.5)
    _, nt = add_text(slide, USABLE_L, 5.95, 11.0, 0.45)
    para(nt, [("说到底：", {"color": ORANGE, "bold": True}),
              ("真实度和交互性，过去很难同时拿到——这正是我们要解决的问题。", {})],
         size=13, color=TEXT, first=True)
    chrome(slide, page, section="01 · 背景")


# ---------------- 背景：SIL / HIL 两个仿真层级 ----------------
def build_silhil_intro(slide, page):
    title(slide, [("仿真的两个层级：SIL 软件在环 vs HIL 硬件在环", {})])
    lead(slide, [("同一套场景，先在纯软件里大规模快跑（SIL），再接上真实域控/芯片验证性能与部署（HIL）——保真度递增、规模递减。", {})])
    fx, fw = 0.72, 5.6
    fh = fw * 9 / 16
    slide.shapes.add_picture(os.path.join(AID, "sil_hil.png"),
                             Inches(fx), Inches(2.05), Inches(fw), Inches(fh))
    add_box(slide, fx, 2.05, fw, fh, line=BORDER, line_w=1)
    _, cap = add_text(slide, fx + 0.1, 2.05 + fh + 0.06, fw - 0.2, 0.3)
    para(cap, [("SIL 纯软件并行（云端）  vs  HIL 接入真实域控/台架", {})],
         size=9.5, color=CYAN, font=MONO, first=True)
    bx = 7.25
    cards = [("SIL · 软件在环（Software-in-the-Loop）",
              "纯软件环境，模型在服务器/云上跑、无真实硬件；快、可大规模并行、成本低——适合批量回归与长里程验证。", CYAN),
             ("HIL · 硬件在环（Hardware-in-the-Loop）",
              "接入真实域控/芯片/台架，硬件在回路中；更接近真车，验证时序、性能与部署——慢、贵、规模受限。", ORANGE),
             ("两级怎么配合",
              "SIL 先大规模筛（广而快），HIL 再高保真验（真而准）；两级互补，共同支撑发版准出。", GREEN)]
    for i, (h, b, c) in enumerate(cards):
        card_block(slide, bx, 2.05 + i * 1.18, USABLE_R - bx, 1.02, h, b,
                   accent=c, head_size=13, body_size=11.2)
    _, nt = add_text(slide, USABLE_L, 5.95, 11.0, 0.45)
    para(nt, [("一句话：", {"color": ORANGE, "bold": True}),
              ("软件在环管“广和快”，硬件在环管“真和准”，两者接力把模型考到位。", {})],
         size=13, color=TEXT, first=True)
    chrome(slide, page, section="01 · 背景")


# ---------------- 技术演进：SimWorld 技术全景 ----------------
def build_simworld_arch(slide, page):
    title(slide, [("G3 · 落到我们：SimWorld 算法框架", {})])
    lead(slide, [("相机 + LiDAR 采集 → 场景图 3DGS 重建 → 扩散修复 → 仿真接口部署；并支持场景编辑——NPC 变轨、资产替换、无限变体。", {})])
    # 左：场景图分解
    panel(slide, USABLE_L, 1.95, 3.5, 4.25, accent=CYAN)
    _, tf = add_text(slide, USABLE_L + 0.22, 2.16, 3.06, 3.9)
    para(tf, [("场景图分解 Scene Graph", {})], size=13, color=CYAN, bold=True,
         font=MONO, first=True, space_after=10)
    for nm, desc in [("背景 Background", "静态建筑/植被"), ("地面 Ground", "2DGS 平坦约束"),
                     ("车辆 RigidNodes", "刚体轨迹建模"), ("行人 SMPL", "人体可变形"),
                     ("天空 Sky", "MLP / CubeMap"), ("交通灯 TLR", "专用高斯")]:
        para(tf, [("· " + nm + "  ", {"color": TEXT, "bold": True}), (desc, {"color": MUTED})],
             size=11.5, space_after=7, line=1.15)
    # 右：四大模块
    bx = 4.55
    mods = [("双 Pipeline", "LiDAR 主力（点云高精）· Vision 纯视觉（MVSNet 深度估计），覆盖有/无激光雷达车型。", CYAN),
            ("三训练模型", "Street Gaussians（优化式）· Reconic/OmniRe（场景图联合训练）· EvoSplat（前馈，免逐场景优化）。", ORANGE),
            ("扩散修复", "Difix3D+ / NVFixer 单步扩散修复渲染伪影；可与 3DGS 联合训练形成“渲染→修复→监督”闭环。", GREEN),
            ("部署落地", "自研 xpeng_raster 光栅化 + 仿真接口（新视角渲染 / 重畸变 / DCCF 平滑），对接 IPS 量产部署。", CYAN)]
    for i, (h, b, c) in enumerate(mods):
        card_block(slide, bx, 1.95 + i * 1.08, USABLE_R - bx, 0.95, h, b,
                   accent=c, head_size=13, body_size=11)
    ref(slide, "Street Gaussians (ECCV’24) · OmniRe · Difix3D+ (NVIDIA CVPR’25) · G3R / EvoSplat 前馈重建")
    chrome(slide, page, section="02 · 技术演进")


# ---------------- 现状：业务为王 ----------------
def build_business(slide, page):
    title(slide, [("业务为王：仿真如何服务量产", {})])
    lead(slide, [("我们做重建不是为了炫技——每一项技术，最终都要对应量产链路上的一个具体业务角色。", {})])
    fx, fw = 0.72, 5.6
    fh = fw * 9 / 16
    slide.shapes.add_picture(os.path.join(AID, "business_loop.png"),
                             Inches(fx), Inches(2.05), Inches(fw), Inches(fh))
    add_box(slide, fx, 2.05, fw, fh, line=BORDER, line_w=1)
    _, cap = add_text(slide, fx + 0.1, 2.05 + fh + 0.06, fw - 0.2, 0.3)
    para(cap, [("路测数据 → 仿真考场 → 版本准出 → 上车，闭环服务量产", {})],
         size=9.5, color=CYAN, font=MONO, first=True)
    bx = 7.25
    roles = [("版本准出 Gating", "模型上车前先在仿真考场跑千万级里程，红线不过不准出。", ORANGE),
             ("长尾挖掘与轰炸", "把路测 Hard Case 泛化成上百变体，专打模型软肋。", CYAN),
             ("归因与定位", "接管率 / 异常自动归因到具体场景，反哺数据闭环。", GREEN),
             ("降本提效", "虚拟里程替代高成本真车路测，迭代周期从“周级”压到“天级”。", ORANGE)]
    for i, (h, b, c) in enumerate(roles):
        card_block(slide, bx, 2.05 + i * 0.92, USABLE_R - bx, 0.8, h, b,
                   accent=c, head_size=12.5, body_size=10.8)
    _, nt = add_text(slide, USABLE_L, 5.95, 11.0, 0.45)
    para(nt, [("简单说：", {"color": ORANGE, "bold": True}),
              ("业务定义“考什么”，技术定义“怎么考”——技术服务业务，这是我们的底层逻辑。", {})],
         size=13, color=TEXT, first=True)
    chrome(slide, page, section="03 · 现状与落地")


# ---------------- 结尾：Q&A ----------------
def build_qa(slide, page):
    slide.shapes.add_picture(os.path.join(AID, "cover_ai.png"),
                             Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    overlay(slide, 0, 0, 13.333, 7.5, rgb=BG, opacity=62)
    _, tt = add_text(slide, 0.85, 2.5, 11.6, 1.4, anchor=MSO_ANCHOR.MIDDLE)
    para(tt, [("Q & A", {})], size=72, color=ORANGE, bold=True,
         align=PP_ALIGN.CENTER, first=True)
    _, ts = add_text(slide, 0.85, 4.2, 11.6, 0.6)
    para(ts, [("谢谢聆听 · 欢迎提问与交流", {})], size=22, color=TEXT, bold=True,
         align=PP_ALIGN.CENTER, first=True)
    _, td = add_text(slide, 0.85, 5.0, 11.6, 0.4)
    para(td, [("XPENG · 自动驾驶仿真算法", {})], size=13, color=MUTED, font=MONO,
         align=PP_ALIGN.CENTER, first=True)
    chrome(slide, page, dark=True)


# ---------------- 五代世代页（G1/G2/G3概念/G5 通用版式）----------------
def _gen_slide(slide, page, gtitle, ltxt, img, cards, axis, labels=None, do_ref=None):
    title(slide, [(gtitle, {})])
    lead(slide, [(ltxt, {})])
    fx, fw = 0.72, 5.6
    fh = fw * 9 / 16
    slide.shapes.add_picture(os.path.join(AID, img),
                             Inches(fx), Inches(2.05), Inches(fw), Inches(fh))
    add_box(slide, fx, 2.05, fw, fh, line=BORDER, line_w=1)
    if labels:
        add_box(slide, fx, 2.05, fw / 2, 0.34, fill=PANEL2)
        add_box(slide, fx + fw / 2, 2.05, fw / 2, 0.34, fill=PANEL2)
        _, l1 = add_text(slide, fx + 0.1, 2.07, fw / 2 - 0.2, 0.3, anchor=MSO_ANCHOR.MIDDLE)
        para(l1, [(labels[0], {})], size=9.5, color=MUTED, font=MONO, first=True)
        _, l2 = add_text(slide, fx + fw / 2 + 0.1, 2.07, fw / 2 - 0.2, 0.3, anchor=MSO_ANCHOR.MIDDLE)
        para(l2, [(labels[1], {})], size=9.5, color=CYAN, font=MONO, first=True)
    bx = 7.25
    for i, (h, b, c) in enumerate(cards):
        card_block(slide, bx, 2.05 + i * 1.18, USABLE_R - bx, 1.02, h, b,
                   accent=c, head_size=13.5, body_size=11.5)
    _, nt = add_text(slide, USABLE_L, 5.98, 11.2, 0.42)
    para(nt, [("双轴定位：", {"color": ORANGE, "bold": True}), (axis, {})],
         size=12.5, color=TEXT, first=True)
    if do_ref:
        ref(slide, do_ref, t=6.5)
    chrome(slide, page, section="02 · 技术演进")


def build_g1(slide, page):
    _gen_slide(slide, page,
        "G1 · 规则仿真：用脚本“写”出一个世界",
        "第一代仿真完全基于规则——工程师手写场景布局与 NPC 行为逻辑。",
        "g1_rules.png",
        [("概念", "纯规则 / 脚本驱动，NPC 走预设的 If-Else 决策。", CYAN),
         ("特点", "完全可控、轻量、可复现；想测什么就写什么。", GREEN),
         ("局限", "无真实感、长尾全靠手写，几乎无法泛化。", ORANGE)],
        "真实度 ▁ 低　·　交互性 ▁ 低（脚本化、非自主）")


def build_g2(slide, page):
    _gen_slide(slide, page,
        "G2 · 游戏引擎：可交互，但“太干净”",
        "第二代用 UE / CARLA 等游戏引擎搭建可交互的 3D 虚拟环境。",
        "g2_engine.png",
        [("概念", "手工搭 3D 资产 + 物理引擎，NPC 可做规则博弈。", CYAN),
         ("特点", "首次实现闭环交互——环境会随本车动作改变。", GREEN),
         ("局限", "Sim2Real gap 大，画面“游戏感”重，视觉大模型不认。", ORANGE)],
        "交互性 ▃↑ 提升　·　真实度 ▁ 仍受限（CGI 失真）")


def build_g3_concept(slide, page):
    _gen_slide(slide, page,
        "G3 · Gauss + Diffusion：真实度的质变（我们所处）",
        "第三代从真实路采用 3DGS 显式重建场景，再用扩散模型修复——既真实，又可编辑、可实时。",
        "nerf_3dgs.png",
        [("概念", "显式 3DGS 重建（替代隐式 NeRF）+ Difix/NVFixer 扩散修复。", CYAN),
         ("特点", "照片级真实 + 100+FPS 实时 + 点级可编辑。", GREEN),
         ("为何质变", "真实度直追实拍，并首次让“真实”与“可编辑交互”兼得。", ORANGE)],
        "真实度 ▆↑↑　·　交互性 ▄↑（点级可编辑）",
        labels=("NeRF 隐式（模糊/慢）", "3DGS 显式（清晰/极速）"),
        do_ref="NeRF (ECCV’20) · 3D Gaussian Splatting (Kerbl, SIGGRAPH’23) · Difix3D+ (NVIDIA CVPR’25)")


def build_g5(slide, page):
    _gen_slide(slide, page,
        "G5 · VLA + RL + Agent：让考场自己进化",
        "第五代——仿真不再只是“被测对象”，而是闭环里能自我迭代的一环。",
        "g5_agent.png",
        [("概念", "VLA 闭环 + 强化学习 + Agent 自动化贯穿研发流。", CYAN),
         ("特点", "AI 自动造考场、数据飞轮自迭代、研发流被重构。", GREEN),
         ("衔接", "这正是第四部分要讲的 VLM + Agent 重塑管线。", ORANGE)],
        "真实度 + 交互性 → 系统级“自进化”（Scaling Law 的工程含义）")


# ====================================================================
#  主装配：按主线顺序构建 25 页
# ====================================================================
def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(s)
    return s


# (builder, kwargs) —— builder 接收 (slide, page) 或 R 内置 (slide)
PLAN = [
    ("cover", build_cover, None),
    ("agenda", build_agenda, None),
    ("div", build_divider, 1),
    ("bg", build_sim_intro, None),
    ("bg", build_silhil_intro, None),
    ("bg", R.s02_paradigm, "01 · 背景"),
    ("bg", build_dataloop_img, None),
    ("bg", build_deadlock_img, None),
    ("div", build_divider, 2),
    ("evo", build_evolution, None),
    ("evo", build_g1, None),
    ("evo", build_g2, None),
    ("evo", build_g3_concept, None),
    ("evo", build_simworld_arch, None),
    ("evo", build_worldmodel_img, None),
    ("evo", build_wm_frontier, None),
    ("evo", build_g5, None),
    ("div", build_divider, 3),
    ("cur", build_team, None),
    ("cur", build_business, None),
    ("cur", build_battle, None),
    ("div", build_divider, 4),
    ("fut", build_vlm_img, None),
    ("fut", build_agents_img, None),
    ("fut", build_agent_future, None),
    ("fut", build_roadmap, None),
    ("fut", R.s14_philosophy, "04 · 未来展望"),
    ("fut", build_mindset, None),
    ("close", build_qa, None),
]


NOTES = {
 1: "【约 0:30】大家好，我是仿真算法的 Li Kun。今天想和大家聊：在自动驾驶进入端到端大模型的时代，我们这群“修考场的人”到底在做什么。一句话主线——我们的使命，是修一座追得上 VLA 模型进化速度的考场。",
 2: "【约 1:00】先看路线图，分四块：第一块背景，为什么旧仿真在端到端时代集体失灵；第二块技术演进，我用仿真的五个世代串起来，一直讲到我们自己的 SimWorld 框架；第三块现状，我们做到哪、怎么服务业务；最后展望未来，AI Agent 怎么让考场自己进化。大家有问题随时打断。",
 3: "【约 0:15】我们先从背景说起——旧考场为什么会失灵。",
 4: "【约 1:30】先花一分钟让非仿真方向的同学也跟上：自动驾驶仿真到底在做什么？一句话，就是在虚拟世界里给模型造一个考场，让它跑里程、撞遍极端场景，把该犯的错在仿真里犯完，而不是拿真车去试——更安全、更便宜、更快。传统有两条路：LogSim 把真车录像原样回放，很真但“死”，环境不会随你的动作变；WorldSim 用游戏引擎手工搭，是“活”的可交互，但画面太假。真实度和交互性，长期就是鱼和熊掌——这正是我们要破的局。",
 5: "【约 1:30】那为什么到了 VLA 2.0 端到端，老办法就不行了？过去模块化时代，感知输出 BBox、车道线，逻辑写在 C++ 规则里，出问题能 Review 代码。但端到端把这些“安全拐杖”全抽走了：驾驶逻辑压进了神经网络权重，看不见也改不动；从感知到决策一步到位，中间没有接口让我们断言“它到底看没看到行人”；行为还会涌现、可能突然失灵，光读代码无法预判。结论就是——静态 Review 失效了，只能靠考场去考。",
 6: "【约 1:15】而且仿真不只是用来“测”。对 VLA 来说，仿真更是它做强化学习的虚拟教练，是模型进化的强核燃料。这是一个数据飞轮：每天在仿真里跑千万级里程给新版本做模拟考，快速判断能不能上车；把路测捞到的疑难场景泛化成上百变体，专挑软肋反复练；最后留下真正变强的、淘汰刷分却变笨的。飞轮转得越快，模型进化越快。",
 7: "【约 1:15】但这里有个技术债危机：传统两套仿真面对 VLA 都是死局。LogSim 开环回放，环境不随车变，而且 VLA 没有 BBox 接口，根本没法 Assert；WorldSim 画面太假，视觉大模型不认这种“游戏画面”。所以我们必须找第三条路——既有 LogSim 的真实度，又有 WorldSim 的交互性。这就引出第二部分。",
 8: "【约 0:15】进入第二部分，技术演进。",
 9: "【约 1:00】我用一条主线串起整个技术演进：仿真的五个世代，每一代都在同时逼近两件事——更高的真实度、更强的交互性。G1 纯规则，G2 游戏引擎，G3 是 Gauss 加 Diffusion 也就是我们现在所处的阶段，G4 世界模型，G5 是 VLA+RL+Agent 闭环自进化。下面一代一代讲。",
 10: "【约 0:45】第一代，规则仿真。最早就是工程师纯手写——场景、NPC 行为全是脚本，走 If-Else。好处是完全可控、轻量、可复现。但缺点致命：没有真实感，长尾全靠手写，几乎无法泛化。真实度、交互性都很低。",
 11: "【约 0:45】第二代，游戏引擎。用 UE、CARLA 搭可交互的 3D 环境，第一次实现闭环——环境会随本车动作改变。交互性上来了，但真实度受限：Sim2Real gap 大，画面“游戏感”重，视觉大模型不买账。",
 12: "【约 1:15】第三代，也是我们现在的主场——Gauss 加 Diffusion。左边图很直观：过去 NeRF 是隐式表达，模糊又慢；现在用 3DGS 显式重建，从真实路采里把场景“焊”出来，清晰又极速，再用 Difix 这类扩散模型修掉伪影。结果是质变：照片级真实、100+FPS 实时、还能点级编辑。第一次，真实度和可编辑的交互性兼得了。",
 13: "【约 1:30】这一代落到我们身上，就是 SimWorld 这套框架。流程：相机加激光雷达采集 → 场景图 3DGS 重建 → 扩散修复 → 仿真接口部署，还支持场景编辑、NPC 变轨、资产替换。几个关键点：把场景拆成背景/地面/车辆/行人/天空/交通灯分别建模；有 LiDAR 和纯视觉两条 pipeline 覆盖所有车型；三个训练模型各有分工；最后自研光栅化器加仿真接口对接量产部署。这是团队真正在写的东西。",
 14: "【约 1:00】第四代，世界模型。前面都是“重建”已有场景，世界模型更进一步，让环境本身有物理常识、会预测未来。通俗讲会三件事：会“脑补”未来——给它当前画面和本车动作，它能推演接下来几秒；会“讲道理”——像真人一样博弈而不是死板 If-Else；会“变天”——一键生成各种天气光照。这是解决“仿真太假、NPC 太呆”的终极武器。",
 15: "【约 1:15】世界模型也是 2026 最热的话题，但业界分两派。理解派以 LeCun 的 AMI Labs 为代表，押注 JEPA、追求真正理解物理；生成派像 DeepMind 的 Genie 3、李飞飞的 World Labs，做可交互的世界生成器。核心争议是：能生成逼真画面≠真正理解物理。这恰好定义了我们的位置——不追 AGI 级理解，做 3D-grounded 生成，几何一致、可编辑、闭环可控，把世界模型做成一座“测得准”的考场。",
 16: "【约 0:45】第五代，VLA+RL+Agent。到这一步，仿真不再只是被测对象，而是闭环里自我进化的发动机：AI 自动造考场，数据飞轮自迭代，研发流被重构。真实度和交互性最终汇成系统级的“自进化”——这就是 Scaling Law 真正的工程含义，也是第四部分要展开的。",
 17: "【约 0:15】第三部分，看看我们现在做到哪、怎么服务业务。",
 18: "【约 1:00】我们是一支在上海的团队，分四条产品线，核心理念是把前沿算法焊死在量产链路上。场景与生产线，把真实路采重建成能直接用的场景；SIL 软件在环，纯软件跑模型回归；HIL 硬件在环，接真实域控做台架仿真；Agents 智能体线，把研发里的重复劳动自动化。分工逻辑：业务 Owner 定义“考什么”，技术 Owner 定义“怎么考”。我们不是辅助工具，是端到端时代的准入裁判。",
 19: "【约 1:15】这里我想强调——业务为王。做重建、做修复不是为了炫技，每项技术都要落到量产链路的一个业务角色上。四个角色：版本准出 Gating，模型上车前先在考场跑千万级里程，红线不过不准出；长尾挖掘，把 Hard Case 泛化成上百变体专打软肋；归因定位，自动把接管率异常归因到具体场景；降本提效，用虚拟里程替代真车路测，迭代从周级压到天级。一句话，业务定义考什么，技术定义怎么考。",
 20: "【约 1:15】光说理念不够，看 Q2 实测战报，每个数字背后都是给量产省下的成本、追回的时间、兜住的风险。单场景重建 91 分钟、提速 2.6 倍，当天建当天用；HIL 仿真 1000 公里每天，替真车路测；重建保真度 PSNR 32 分贝以上，接近实拍可直接训模型；坏场景质检自动召回 90%；线上问题闭环复现率 74%，正冲 80%；12 个 Agent 专项已上线 4 个。",
 21: "【约 0:15】最后，展望未来——怎么让考场自己进化。",
 22: "【约 1:00】未来核心是 AI 重塑整个研发管线。第一步，VLM 自动化资产沉淀。痛点是人力看不过来每天 TB 级的路测 Log。VLM 可以自动扫描、输出“无保护左转”这类高价值语义标签；自动找出 VLA 表现差的薄弱片段直接送重训；把海量 Log 精炼成高优先级训练集——把“人看不过来”变成“机器筛得过来”。",
 23: "【约 1:15】第二步是 Agent。其实 Agent 今天已经在替我们干活了：Coding Agent 做代码库改造和工作流自动化；Case 打标 Agent 自动给场景打标签；复现率分析 Agent 判断线上问题能否在仿真复现；Metric 分析 Agent 自动读大盘指标做归因；Smart Planner 用自然语言编排仿真任务；还有 Text2SQL 的 Data Agent，用大白话就能查大盘。这些不是设想，是在用的。",
 24: "【约 1:00】往前看，Agent 还能做什么？我的判断是它会从“帮手”进化成一支自主的工程军团：从“发现问题”升级到“修好问题”，自己定位根因、提修复 MR，人只 Review；多个 Agent 协同，Planner、Coder、Reviewer、Tester 分工端到端跑完一个需求；甚至能按模型当前弱点，自动出由易到难的“考题课程”驱动模型自进化；结合世界模型，一句话造出符合物理的长尾场景。终极形态，是我们从“修考场”升级为“指挥一支 AI 工程军团”。",
 25: "【约 1:00】把这些串成路线图，三步走。NOW 夯实考场：复现率冲 80%、SIL 效率到 1:25、HIL 稳定、车型全覆盖。NEXT 智能化管线：Agent 自动复现归因修复、VLM 资产沉淀闭环、12 专项全量上线。FUTURE 生成式世界：世界模型生成长尾、DiT 加 3DGS 混合落地，仿真变成一座“数据工厂”。",
 26: "【约 0:45】讲点工程哲学。有句话我很喜欢：算法团队负责跑分，我们负责修考场。我们的使命，是用最前沿的 3DGS 和世界模型，在充满物理噪声的平行世界里，为自动驾驶打磨出最“确定”的准出标准。两条原则：接口第一，灵活性来自极简清晰的抽象；不迷信日志，在闭环数据和真实物理面前，任何主观借口都苍白。",
 27: "【约 1:00】最后送给在座硕博同学三条箴言，关于怎么在“不确定”里做出“确定”。第一，不确定到确定：大模型本质是概率，工程本质是可复现，用闭环和裁判机制把“大概率对”变成“可验证的确定”。第二，论文到产线：能跑 demo 不算赢，能在量产线稳定出数字才算。第三，单点到闭环：别迷恋单个算法的 SOTA，要建会自己进化的系统——这才是 Scaling Law 真正的工程含义。",
 28: "【Q&A · 余下时间】我的分享就到这里，谢谢大家！接下来是问答时间，欢迎就仿真、3DGS、世界模型、Agent 任何话题交流。",
}


def main():
    prs = Presentation()
    prs.slide_width = Emu(12191695)
    prs.slide_height = Emu(6858000)
    NEW = {build_cover, build_agenda, build_divider, build_evolution, build_team,
           build_battle, build_roadmap, build_mindset, build_pipeline_img,
           build_worldmodel_img, build_dataloop_img, build_nerf_img,
           build_drift_img, build_vlm_img, build_agents_img,
           build_deadlock_img, build_spatio_img,
           build_wm_frontier, build_wm_grounded,
           build_sim_intro, build_simworld_arch, build_business, build_qa,
           build_silhil_intro,
           build_g1, build_g2, build_g3_concept, build_g5, build_agent_future}
    for i, (kind, builder, arg) in enumerate(PLAN, 1):
        s = new_slide(prs)
        if builder is build_divider:
            builder(s, arg, i)
        elif builder in NEW:
            builder(s, i)
        else:
            builder(s)                       # R 内置内容（HUD 已屏蔽）
            chrome(s, i, section=arg)         # 统一 chrome + logo
        print(f"  p{i:>2} [{kind}] {builder.__name__}")
        if i in NOTES:
            s.notes_slide.notes_text_frame.text = NOTES[i]
    prs.save(OUT)
    print("saved ->", OUT, "| slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
